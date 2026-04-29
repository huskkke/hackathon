#!/usr/bin/env python3
"""
PresentAI — AI-генератор презентаций
Hackathon: Амурский Код 2026 | Кейс: Ростелеком
"""

import asyncio
import base64
import copy
import io
import json
import os
import re
import time
import traceback
import uuid
import zipfile
from pathlib import Path
from typing import Any, Callable, Optional
from xml.sax.saxutils import escape

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

# ── Настройка приложения ──────────────────────────────────────────────────────

SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
WORK_DIR = Path(os.getenv("PRESENTAI_WORK_DIR") or PROJECT_DIR / "temp_files")
WORK_DIR.mkdir(parents=True, exist_ok=True)

# RT API используется только при наличии токена: без него приложение уходит в локальный запасной режим.
RT_API_BASE = "https://ai.rt.ru/api/1.0"
RT_LLM_MODEL = "Qwen/Qwen2.5-72B-Instruct"
# У RT LLM есть чувствительность к большим payload: лимиты ниже держат запросы
# в стабильном диапазоне и уменьшают шанс HTTP 400 на длинных документах.
RT_LLM_MAX_NEW_TOKENS = 1536
RT_LLM_MAX_PROMPT_CHARS = 24000


def _env_int(name: str, default: int, minimum: int, maximum: int) -> int:
    try:
        value = int(os.getenv(name, str(default)))
    except Exception:
        value = default
    return max(minimum, min(maximum, value))


# Генерация изображений потенциально самая долгая часть пайплайна. Эти лимиты
# не дают Yandex ART или /download заблокировать всю презентацию на много минут.
RT_IMAGE_POST_TIMEOUT = _env_int("RT_IMAGE_POST_TIMEOUT", 20, 5, 90)
RT_IMAGE_REQUEST_TIMEOUT = _env_int("RT_IMAGE_REQUEST_TIMEOUT", 12, 3, 60)
RT_IMAGE_DOWNLOAD_REQUEST_TIMEOUT = _env_int("RT_IMAGE_DOWNLOAD_REQUEST_TIMEOUT", 20, 5, 60)
RT_IMAGE_DOWNLOAD_TIMEOUT = _env_int("RT_IMAGE_DOWNLOAD_TIMEOUT", 60, 15, 180)
RT_IMAGE_SINGLE_TIMEOUT = _env_int("RT_IMAGE_SINGLE_TIMEOUT", 70, 20, 180)
RT_IMAGE_BATCH_TIMEOUT = _env_int("RT_IMAGE_BATCH_TIMEOUT", 270, 60, 600)
RT_IMAGE_MAX_IMAGES = _env_int("RT_IMAGE_MAX_IMAGES", 20, 1, 50)
# Fallback на SD включён по умолчанию: если yaArt недоступен (rate limit/таймаут), пробуем SD.
RT_IMAGE_ALLOW_BACKEND_FALLBACK = os.getenv("RT_IMAGE_ALLOW_BACKEND_FALLBACK", "1") == "1"
# При 429 (rate limit) делаем паузу и повторяем попытку в рамках общего дедлайна.
RT_IMAGE_RATELIMIT_PAUSE = _env_int("RT_IMAGE_RATELIMIT_PAUSE", 8, 1, 60)
RT_IMAGE_POST_RETRIES = _env_int("RT_IMAGE_POST_RETRIES", 1, 0, 3)
# Circuit breaker: сколько подряд ошибок одного бэкенда разрешено до его отключения в батче.
RT_IMAGE_CIRCUIT_BREAKER_THRESHOLD = _env_int("RT_IMAGE_CIRCUIT_BREAKER_THRESHOLD", 2, 1, 10)

app = FastAPI(title="PresentAI")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

JOBS: dict[str, dict] = {}
# Токены и runtime-контекст живут только в памяти процесса и не пишутся в JSON/PPTX.
SESSION_SECRETS: dict[str, dict] = {}


ProgressCallback = Callable[[int, str, str], None]


def _job_update(
    job_id: str,
    progress: Optional[int] = None,
    title: Optional[str] = None,
    detail: Optional[str] = None,
    status: Optional[str] = None,
    result: Optional[dict] = None,
    error: Optional[str] = None,
) -> None:
    """Обновляет публичный статус фоновой генерации без раскрытия токенов и файлов."""
    job = JOBS.get(job_id)
    if not job:
        return
    if progress is not None:
        current = int(job.get("progress", 0) or 0)
        job["progress"] = max(current, max(0, min(100, int(progress))))
    if title is not None:
        job["title"] = title
    if detail is not None:
        job["detail"] = detail
    if status is not None:
        job["status"] = status
    if result is not None:
        job["result"] = result
    if error is not None:
        job["error"] = error
    job["updated_at"] = time.time()


def _job_public_status(job_id: str, job: dict) -> dict:
    """Возвращает только данные, которые безопасно показывать странице загрузки."""
    result = job.get("result") if isinstance(job.get("result"), dict) else {}
    return {
        "job_id": job_id,
        "status": job.get("status", "queued"),
        "progress": int(job.get("progress", 0) or 0),
        "title": job.get("title", "Ожидаем запуска"),
        "detail": job.get("detail", ""),
        "error": job.get("error", ""),
        "session_id": result.get("session_id", ""),
        "created_at": job.get("created_at"),
        "updated_at": job.get("updated_at"),
    }

# ── Извлечение текста из документов ───────────────────────────────────────────

def extract_text_from_pdf(data: bytes) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        texts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
        return "\n\n".join(texts)[:8000]
    except Exception as e:
        return f"[PDF извлечение не удалось: {e}]"

def extract_text_from_docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras)[:8000]
    except Exception as e:
        return f"[DOCX извлечение не удалось: {e}]"

# ── LLM: генерация структуры слайдов ─────────────────────────────────────────

SYSTEM_PROMPT = """You are an expert presentation designer and content strategist.
You generate structured JSON for professional presentations. Always respond ONLY with valid JSON.
No markdown, no explanations, just the JSON object."""

DEFAULT_STYLE = "deep_neon"
# Старый ключ оставлен как алиас, чтобы ранее созданные сессии и ссылки не ломались.
STYLE_ALIASES = {
    "rostelecom": DEFAULT_STYLE,
}
STYLE_LABELS = {
    DEFAULT_STYLE: "Глубокий неон",
    "modern": "Тёмная волна",
    "corporate": "Деловая",
    "minimal": "Минимализм",
    "tech": "Технологичная",
    "creative": "Креативная",
}


def normalize_style(style: str) -> str:
    value = (style or DEFAULT_STYLE).strip()
    return STYLE_ALIASES.get(value, value)


def style_label(style: str) -> str:
    normalized = normalize_style(style)
    return STYLE_LABELS.get(normalized, STYLE_LABELS[DEFAULT_STYLE])


def build_generation_prompt(user_prompt: str, doc_text: str, slide_count: int,
                             style: str, tone: str) -> str:
    tone_desc = {
        "professional": "профессиональный, деловой, чёткий",
        "creative": "творческий, вдохновляющий, образный",
        "academic": "академический, точный, аналитический",
        "casual": "дружелюбный, понятный, разговорный",
        "persuasive": "убедительный, мотивирующий, динамичный",
    }.get(tone, "профессиональный")

    doc_section = ""
    if doc_text.strip():
        doc_section = f"\n\nДОКУМЕНТ (источник контента):\n---\n{doc_text[:6000]}\n---"

    return f"""Верни только валидный JSON без markdown.
Создай презентацию на русском языке.
Тема: {user_prompt}
Количество слайдов: {slide_count}
Стиль: {style_label(style)}
Тон: {tone_desc}{doc_section}

Формат:
{{
  "presentation_title": "...",
  "slides": [
    {{"layout":"title","title":"...","subtitle":"..."}},
    {{"layout":"content","title":"...","bullets":["...","...","..."],"image_prompt":"english image prompt"}},
    {{"layout":"two_column","title":"...","leftTitle":"...","leftContent":["...","..."],"rightTitle":"...","rightContent":["...","..."],"image_prompt":"english image prompt"}},
    {{"layout":"stats","title":"...","stats":[{{"value":"...","label":"..."}}],"content":"...","image_prompt":"english image prompt"}},
    {{"layout":"conclusion","title":"...","content":"..."}}
  ]
}}

Правила: ровно {slide_count} слайдов; первый layout title; последний layout conclusion; текст слайдов строго по теме; bullets короткие; image_prompt на английском."""

def _extract_json_object(text: str) -> dict:
    text = (text or "").strip()
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text).strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        start = text.find("{")
        end = text.rfind("}")
        if start >= 0 and end > start:
            return json.loads(text[start:end + 1])
        raise

def _find_text_in_response(data: object) -> str:
    if isinstance(data, str):
        return data
    if isinstance(data, list):
        parts = [_find_text_in_response(item) for item in data]
        return "\n".join(part for part in parts if part)
    if isinstance(data, dict):
        message = data.get("message")
        if isinstance(message, dict):
            content = message.get("content")
            if isinstance(content, str) and content.strip():
                return content
        for key in ("content", "text", "answer", "response", "message", "generated_text"):
            value = data.get(key)
            if isinstance(value, str) and value.strip():
                return value
        for value in data.values():
            found = _find_text_in_response(value)
            if found:
                return found
    return ""

def _decode_rt_response_body(resp: object, limit: int = 900) -> str:
    """Декодирует ошибки RT API как UTF-8, чтобы русские сообщения были читаемыми."""
    if resp is None:
        return ""
    try:
        content = getattr(resp, "content", b"") or b""
        if isinstance(content, bytes) and content:
            text = content.decode("utf-8", errors="replace")
        else:
            text = str(getattr(resp, "text", "") or "")
    except Exception:
        text = str(getattr(resp, "text", "") or "")
    return text.strip()[:limit]

def _safe_rt_max_tokens(value: int) -> int:
    """Удерживает max_new_tokens в рабочем диапазоне RT endpoint."""
    try:
        requested = int(value)
    except Exception:
        requested = RT_LLM_MAX_NEW_TOKENS
    return max(64, min(RT_LLM_MAX_NEW_TOKENS, requested))

def _compact_rt_prompt(prompt: object, limit: int = RT_LLM_MAX_PROMPT_CHARS) -> str:
    """Сжимает слишком большой промпт, сохраняя начало и конец контекста."""
    text = str(prompt or "").replace("\x00", "").strip()
    if len(text) <= limit:
        return text
    head = max(1000, limit // 2)
    tail = max(1000, limit - head - 200)
    return text[:head].rstrip() + "\n\n...[контекст сокращён, чтобы запрос прошёл лимит RT API]...\n\n" + text[-tail:].lstrip()

def _call_rt_llm(
    prompt: str,
    token: str,
    system_prompt: str = SYSTEM_PROMPT,
    max_new_tokens: int = 1536,
    temperature: float = 0.25,
) -> str:
    import requests as req

    safe_prompt = _compact_rt_prompt(prompt)
    safe_max_tokens = _safe_rt_max_tokens(max_new_tokens)
    content_item = {"type": "text", "text": safe_prompt}

    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": RT_LLM_MODEL,
            "user_message": safe_prompt,
            "contents": [content_item],
            "system_prompt": system_prompt,
            "max_new_tokens": safe_max_tokens,
            "temperature": temperature,
            "top_k": 40,
            "top_p": 0.9,
            "no_repeat_ngram_size": 15,
            "repetition_penalty": 1.05,
        },
    }
    try:
        resp = req.post(
            f"{RT_API_BASE}/llama/chat",
            headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json; charset=utf-8"},
            json=payload,
            timeout=75,
        )
        resp.raise_for_status()
    except req.HTTPError as e:
        response = e.response
        status = response.status_code if response is not None else ""
        body = _decode_rt_response_body(response)
        raise RuntimeError(f"HTTP {status}: {body}") from e

    raw = _decode_rt_response_body(resp, limit=200000)
    if not raw:
        raise RuntimeError(f"RT API вернул пустой ответ, HTTP {resp.status_code}")
    try:
        data = resp.json()
    except ValueError:
        return raw

    text = _find_text_in_response(data).strip()
    if not text:
        raise RuntimeError(f"RT API ответил без текста. Фрагмент ответа: {raw[:500]}")
    return text

def _call_anthropic(prompt: str) -> str:
    import anthropic

    client = anthropic.Anthropic()
    message = client.messages.create(
        model=os.getenv("ANTHROPIC_MODEL", "claude-sonnet-4-20250514"),
        max_tokens=4096,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}],
    )
    return message.content[0].text

def _fallback_slide_structure(user_prompt: str, doc_text: str, slide_count: int,
                              style: str, tone: str) -> dict:
    source = re.sub(r"\s+", " ", doc_text).strip()
    topic = user_prompt.strip(" .") or "выбранная тема"
    facts = [item.strip(" .") for item in re.split(r"[.\n;]+", source) if len(item.strip()) > 30]
    if not facts:
        facts = [
            f"Тема требует краткого обзора: что происходит, почему это важно и кого затрагивает",
            f"Для раскрытия темы «{topic}» важно показать причины, последствия и текущие вызовы",
            f"Отдельный блок стоит посвятить статистике, географии и динамике изменений",
            f"Практическая часть должна включать меры профилактики, реагирования и снижения рисков",
        ]

    title = user_prompt[:90].strip(" .") or "Презентация"
    slides = [{
        "layout": "title",
        "title": title,
        "subtitle": f"Краткий аналитический обзор по теме",
        "image_prompt": None,
    }]

    middle_count = max(1, slide_count - 2)
    for i in range(middle_count):
        fact = facts[i % len(facts)]
        if i % 4 == 1:
            slides.append({
                "layout": "two_column",
                "title": "Причины и последствия",
                "leftTitle": "Ключевые причины",
                "leftContent": ["Природные и климатические факторы", "Человеческий фактор", "Недостаточная профилактика"],
                "rightTitle": "Последствия",
                "rightContent": ["Риски для людей и инфраструктуры", "Экономический ущерб", "Экологические потери"],
                "image_prompt": None,
            })
        elif i % 4 == 2:
            slides.append({
                "layout": "stats",
                "title": "Что важно показать цифрами",
                "stats": [
                    {"value": "1", "label": "Масштаб проблемы"},
                    {"value": "2", "label": "Затронутые регионы"},
                    {"value": "3", "label": "Динамика и прогноз"},
                ],
                "content": fact,
                "image_prompt": None,
            })
        elif i % 4 == 3:
            slides.append({
                "layout": "quote",
                "title": "Ключевая идея",
                "quote": fact,
                "image_prompt": None,
            })
        else:
            slides.append({
                "layout": "content",
                "title": f"Слайд {i + 2}",
                "bullets": [
                    fact[:110],
                    f"Свяжите тезис с темой «{topic}» через конкретные примеры",
                    "Добавьте вывод: что нужно сделать дальше",
                ],
                "image_prompt": f"realistic editorial illustration about {topic}, 16:9 presentation slide",
            })

    slides.append({
        "layout": "conclusion",
        "title": "Выводы",
        "content": f"Тема «{topic}» требует системного анализа, профилактики и понятного плана действий.",
        "image_prompt": None,
    })
    return {"presentation_title": title, "metadata": {"title": title, "author": "PresentAI", "contact": ""}, "slides": slides[:slide_count]}

def _structure_from_plain_text(text: str, user_prompt: str, slide_count: int, style: str, tone: str) -> dict:
    chunks = [part.strip(" -\n\t") for part in re.split(r"\n+|(?<=[.!?])\s+", text) if len(part.strip()) > 20]
    if not chunks:
        return _fallback_slide_structure(user_prompt, text, slide_count, style, tone)
    structure = _fallback_slide_structure(user_prompt, text, slide_count, style, tone)
    slides = structure["slides"]
    for idx, chunk in enumerate(chunks[:max(1, slide_count - 2)], start=1):
        if idx >= len(slides) - 1:
            break
        slides[idx] = {
            "layout": "content",
            "title": f"Ключевой тезис {idx}",
            "bullets": [chunk[:150]],
            "image_prompt": f"editorial presentation image about {user_prompt}, 16:9",
        }
    structure["slides"] = slides
    return structure

def _normalize_slide_structure(data: dict, slide_count: int) -> dict:
    if not isinstance(data, dict):
        raise ValueError("LLM вернула не JSON-объект")
    slides = data.get("slides")
    if not isinstance(slides, list):
        raise ValueError("В JSON нет массива slides")
    if len(slides) > slide_count:
        conclusion = next((s for s in reversed(slides) if s.get("layout") == "conclusion"), slides[-1])
        slides = slides[:max(1, slide_count - 1)] + [conclusion]
    while len(slides) < slide_count:
        slides.append({"layout": "content", "title": f"Слайд {len(slides) + 1}", "bullets": [], "image_prompt": None})
    if slides:
        slides[0]["layout"] = "title"
        slides[-1]["layout"] = "conclusion"
    data["slides"] = slides
    data.setdefault("presentation_title", slides[0].get("title", "Презентация") if slides else "Презентация")
    data.setdefault("metadata", {"title": data["presentation_title"], "author": "PresentAI", "contact": ""})
    return data


# ── Второй LLM-контролёр качества и редактор слайдов ─────────────────────────

QUALITY_SYSTEM_PROMPT = """Ты — независимый второй LLM-агент контроля качества презентаций.
Твоя задача — проверять структуру, достоверность, читаемость, визуальный баланс и при необходимости редактировать JSON слайдов.
Всегда отвечай только валидным JSON без markdown, пояснений и лишнего текста."""

SINGLE_SLIDE_EDITOR_SYSTEM_PROMPT = """Ты — LLM-редактор одного слайда презентации.
Редактируй только указанный слайд, не меняй всю презентацию. Возвращай только валидный JSON без markdown."""

ALLOWED_LAYOUTS = {"title", "content", "two_column", "stats", "quote", "section_break", "conclusion", "section", "image"}

SLIDE_TEXT_LIMITS = {
    "title": {"title": 78, "body": 155, "items": 1, "item": 155},
    "content": {"title": 82, "body": 520, "items": 5, "item": 115},
    "two_column": {"title": 78, "body": 620, "items": 4, "item": 95},
    "stats": {"title": 78, "body": 360, "items": 3, "item": 70},
    "quote": {"title": 78, "body": 360, "items": 1, "item": 360},
    "section_break": {"title": 78, "body": 220, "items": 1, "item": 220},
    "conclusion": {"title": 78, "body": 430, "items": 1, "item": 430},
    "section": {"title": 78, "body": 220, "items": 1, "item": 220},
    "image": {"title": 78, "body": 260, "items": 3, "item": 90},
}

def _clean_text(value: object) -> str:
    value = str(value or "")
    value = re.sub(r"\s+", " ", value).strip()
    return value

def _clip_text(value: object, limit: int, suffix: str = "…") -> str:
    text = _clean_text(value)
    if len(text) <= limit:
        return text
    cut = max(0, limit - len(suffix))
    part = text[:cut].rstrip()
    sentence_cut = max(part.rfind(". "), part.rfind("; "), part.rfind(", "), part.rfind(" — "))
    if sentence_cut >= max(30, int(limit * 0.55)):
        part = part[:sentence_cut].rstrip(" .;,—-")
    return part.rstrip(" .;,—-") + suffix

def _as_text_list(value: object) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        result = []
        for item in value:
            if isinstance(item, dict):
                text = item.get("text") or item.get("label") or item.get("title") or ""
            else:
                text = item
            cleaned = _clean_text(text)
            if cleaned:
                result.append(cleaned)
        return result
    if isinstance(value, str):
        lines = [x.strip(" •-\t") for x in re.split(r"\n+|(?:^|\s)[•▪▫◦‣]\s*", value) if x.strip(" •-\t")]
        if len(lines) <= 1:
            lines = [x.strip(" •-\t") for x in re.split(r";\s+|(?<=[.!?])\s+(?=[А-ЯA-Z0-9])", value) if x.strip(" •-\t")]
        return [_clean_text(x) for x in lines if _clean_text(x)]
    cleaned = _clean_text(value)
    return [cleaned] if cleaned else []

def _slide_plain_text(slide: dict) -> str:
    parts: list[str] = []
    for key in ("title", "subtitle", "content", "quote", "leftTitle", "rightTitle"):
        if slide.get(key):
            parts.append(str(slide.get(key)))
    for key in ("bullets", "leftContent", "rightContent"):
        parts.extend(_as_text_list(slide.get(key)))
    for stat in slide.get("stats", []) or []:
        if isinstance(stat, dict):
            parts.append(str(stat.get("value", "")))
            parts.append(str(stat.get("label", "")))
    return _clean_text(" ".join(parts))

def _slide_density(slide: dict) -> dict:
    text = _slide_plain_text(slide)
    layout = str(slide.get("layout") or "content")
    limits = SLIDE_TEXT_LIMITS.get(layout, SLIDE_TEXT_LIMITS["content"])
    line_count = max(1, text.count("\n") + len(re.findall(r"[.!?;]\s", text)) + 1)
    item_count = 0
    if layout == "two_column":
        item_count = len(_as_text_list(slide.get("leftContent"))) + len(_as_text_list(slide.get("rightContent")))
    elif layout == "stats":
        item_count = len(slide.get("stats") or [])
    elif layout == "content":
        item_count = len(_as_text_list(slide.get("bullets")))
    ratio = max(
        len(text) / max(1, limits["body"]),
        item_count / max(1, limits["items"]),
        line_count / 12,
    )
    return {"chars": len(text), "line_count": line_count, "item_count": item_count, "ratio": round(ratio, 2)}

def _doc_keyword_set(doc_text: str) -> set[str]:
    words = re.findall(r"[A-Za-zА-Яа-яЁё0-9]{4,}", (doc_text or "").lower())
    stop = {
        "которые", "которая", "который", "также", "может", "можно", "нужно",
        "этого", "этой", "если", "были", "будет", "через", "после", "перед",
        "their", "with", "from", "this", "that", "have", "will", "about",
    }
    return {w for w in words if w not in stop}

def _slide_numbers(text: str) -> set[str]:
    return {m.group(0) for m in re.finditer(r"(?<![\w])\d{2,}(?:[.,]\d+)?%?(?![\w])", text or "")}

def _append_issue(issues: list[dict], slide_num: int, severity: str, issue: str, fix: str = "") -> None:
    issues.append({
        "slide": slide_num,
        "severity": severity,
        "issue": issue,
        "fix": fix,
    })

def _normalize_layout_name(layout: object) -> str:
    layout = str(layout or "content").strip().lower().replace("-", "_").replace(" ", "_")
    return layout if layout in ALLOWED_LAYOUTS else "content"

def _sanitize_slide_for_quality(
    slide: dict,
    index: int,
    total: int,
    doc_keywords: Optional[set[str]] = None,
    doc_text: str = "",
) -> tuple[dict, list[dict], list[str]]:
    """Приводит слайд к безопасному виду для PPTX и собирает QA-заметки."""
    s = copy.deepcopy(slide if isinstance(slide, dict) else {})
    issues: list[dict] = []
    corrections: list[str] = []
    slide_num = index + 1

    layout = _normalize_layout_name(s.get("layout"))
    if layout == "section":
        layout = "section_break"
    if index == 0 and layout != "title":
        _append_issue(issues, slide_num, "medium", "Первый слайд был не титульным.", "layout изменён на title")
        corrections.append(f"Слайд {slide_num}: установлен титульный layout.")
        layout = "title"
    if index == total - 1 and layout != "conclusion":
        _append_issue(issues, slide_num, "medium", "Последний слайд был не заключительным.", "layout изменён на conclusion")
        corrections.append(f"Слайд {slide_num}: установлен заключительный layout.")
        layout = "conclusion"
    s["layout"] = layout

    limits = SLIDE_TEXT_LIMITS.get(layout, SLIDE_TEXT_LIMITS["content"])

    title = _clean_text(s.get("title")) or ("Презентация" if layout == "title" else f"Слайд {slide_num}")
    clipped_title = _clip_text(title, limits["title"])
    if clipped_title != title:
        _append_issue(issues, slide_num, "low", "Заголовок слишком длинный для слайда.", "заголовок сокращён")
        corrections.append(f"Слайд {slide_num}: сокращён заголовок.")
    s["title"] = clipped_title

    if layout == "title":
        subtitle = _clean_text(s.get("subtitle") or s.get("content") or "")
        s["subtitle"] = _clip_text(subtitle, limits["body"])
        if len(subtitle) > limits["body"]:
            corrections.append(f"Слайд {slide_num}: сокращён подзаголовок титульного слайда.")
    elif layout == "two_column":
        left_title = _clip_text(s.get("leftTitle") or "Ключевой блок", 42)
        right_title = _clip_text(s.get("rightTitle") or "Детали", 42)
        left = _as_text_list(s.get("leftContent"))
        right = _as_text_list(s.get("rightContent"))
        if not left and s.get("content"):
            items = _as_text_list(s.get("content"))
            half = max(1, (len(items) + 1) // 2)
            left, right = items[:half], items[half:]
        max_each = limits["items"]
        old_count = len(left) + len(right)
        left = [_clip_text(x, limits["item"]) for x in left[:max_each]]
        right = [_clip_text(x, limits["item"]) for x in right[:max_each]]
        if len(left) + len(right) < old_count:
            _append_issue(issues, slide_num, "medium", "В двухколоночном слайде было слишком много пунктов.", "лишние пункты убраны")
            corrections.append(f"Слайд {slide_num}: уменьшено количество пунктов в колонках.")
        s["leftTitle"] = left_title
        s["rightTitle"] = right_title
        s["leftContent"] = left or ["Ключевой тезис"]
        s["rightContent"] = right or ["Практический вывод"]
    elif layout == "stats":
        raw_stats = s.get("stats") if isinstance(s.get("stats"), list) else []
        stats = []
        for n, stat in enumerate(raw_stats[:limits["items"]], start=1):
            if isinstance(stat, dict):
                value = _clip_text(stat.get("value") or str(n), 18)
                label = _clip_text(stat.get("label") or "Показатель", 62)
            else:
                value = str(n)
                label = _clip_text(stat, 62)
            stats.append({"value": value, "label": label})
        if not stats:
            stats = [{"value": "1", "label": "Главный показатель"}, {"value": "2", "label": "Ключевой риск"}, {"value": "3", "label": "Следующий шаг"}]
        if len(raw_stats) > limits["items"]:
            _append_issue(issues, slide_num, "low", "Статистический слайд содержал больше трёх карточек.", "оставлены три ключевых показателя")
            corrections.append(f"Слайд {slide_num}: статистика сокращена до трёх карточек.")
        content = _clean_text(s.get("content") or s.get("subtitle") or "")
        s["stats"] = stats
        s["content"] = _clip_text(content, limits["body"])
    elif layout == "quote":
        quote = _clean_text(s.get("quote") or s.get("content") or s.get("subtitle") or "")
        if not quote:
            quote = "Ключевая мысль презентации должна быть понятной и проверяемой."
        s["quote"] = _clip_text(quote, limits["body"])
    elif layout in {"section_break", "conclusion"}:
        content = _clean_text(s.get("content") or s.get("subtitle") or "")
        s["content"] = _clip_text(content, limits["body"])
        s["subtitle"] = s["content"]
        if layout == "conclusion" and not s["content"]:
            s["content"] = "Суммируйте ключевые выводы, риски и следующий практический шаг."
            s["subtitle"] = s["content"]
    else:
        bullets = _as_text_list(s.get("bullets"))
        if not bullets and s.get("content"):
            bullets = _as_text_list(s.get("content"))
        if not bullets:
            bullets = ["Ключевой тезис", "Практическое значение", "Следующий шаг"]
        old_count = len(bullets)
        bullets = [_clip_text(x, limits["item"]) for x in bullets[:limits["items"]]]
        if old_count > limits["items"]:
            _append_issue(issues, slide_num, "medium", "Слайд содержал слишком много пунктов и мог выглядеть перегруженным.", "количество пунктов сокращено")
            corrections.append(f"Слайд {slide_num}: сокращён список пунктов.")
        s["bullets"] = bullets

    density = _slide_density(s)
    if density["ratio"] > 1.15:
        _append_issue(
            issues,
            slide_num,
            "medium",
            f"Слайд может быть визуально перегружен: {density['chars']} символов.",
            "текст сокращён локальным контролёром",
        )

    if doc_keywords:
        plain = _slide_plain_text(s).lower()
        words = {w for w in re.findall(r"[A-Za-zА-Яа-яЁё0-9]{4,}", plain) if len(w) >= 4}
        meaningful = {w for w in words if not w.isdigit()}
        overlap = meaningful & doc_keywords
        if len(meaningful) >= 8 and len(overlap) <= 1 and index not in (0, total - 1):
            _append_issue(
                issues,
                slide_num,
                "low",
                "Часть текста слайда слабо связана с загруженным документом.",
                "рекомендуется уточнить источники или переформулировать тезисы",
            )
        doc_numbers = _slide_numbers(doc_text)
        for num in sorted(_slide_numbers(plain)):
            if num not in doc_numbers:
                _append_issue(
                    issues,
                    slide_num,
                    "medium",
                    f"Числовое значение «{num}» не найдено в исходном документе.",
                    "число оставлено, но отмечено для проверки пользователем",
                )

    return s, issues, corrections

def _local_quality_review_and_edit(
    structure: dict,
    user_prompt: str,
    doc_text: str,
    slide_count: int,
    style: str,
    tone: str,
) -> dict:
    """Локальный QA-фильтр на случай, если внешний LLM-контролёр недоступен."""
    reviewed = _normalize_slide_structure(copy.deepcopy(structure), slide_count)
    doc_keywords = _doc_keyword_set(doc_text)
    issues: list[dict] = []
    corrections: list[str] = []
    cleaned_slides: list[dict] = []
    for idx, slide in enumerate(reviewed.get("slides", [])):
        cleaned, slide_issues, slide_corrections = _sanitize_slide_for_quality(
            slide, idx, slide_count, doc_keywords, doc_text
        )
        cleaned_slides.append(cleaned)
        issues.extend(slide_issues)
        corrections.extend(slide_corrections)

    reviewed["slides"] = cleaned_slides
    reviewed["presentation_title"] = reviewed.get("presentation_title") or (cleaned_slides[0].get("title") if cleaned_slides else "Презентация")
    reviewed.setdefault("metadata", {"title": reviewed["presentation_title"], "author": "PresentAI", "contact": ""})

    penalty = 0
    for issue in issues:
        penalty += {"high": 18, "medium": 9, "low": 4}.get(issue.get("severity"), 5)
    score = max(55, min(100, 100 - penalty))
    reviewed["quality_review"] = {
        "source": "Локальный QA-контроль",
        "score": score,
        "summary": "Проверены порядок слайдов, плотность текста, базовая связность с источником и потенциальные переполнения.",
        "issues": issues[:40],
        "corrections": corrections[:40],
    }
    reviewed["review_source"] = "Локальный QA-контроль"
    return reviewed

LLM_STRIPPED_RUNTIME_KEYS = {
    "imageData", "image_data", "imageBase64", "base64", "thumbnail", "previewImage",
    "pptx", "file", "blob", "binary", "bytes",
}

def _strip_runtime_fields_for_llm(value: object, key: str = "") -> object:
    """Удаляет base64, изображения и runtime-поля перед отправкой JSON во вторую LLM."""
    if key in LLM_STRIPPED_RUNTIME_KEYS:
        return None
    if isinstance(value, dict):
        cleaned: dict[str, object] = {}
        for k, v in value.items():
            if k in LLM_STRIPPED_RUNTIME_KEYS:
                continue
            cleaned_v = _strip_runtime_fields_for_llm(v, k)
            if cleaned_v is not None:
                cleaned[k] = cleaned_v
        return cleaned
    if isinstance(value, list):
        return [_strip_runtime_fields_for_llm(item, key) for item in value]
    if isinstance(value, str):
        value_str = value.strip()
        lower_key = key.lower()
        if value_str.startswith("image/") or value_str.startswith("data:image"):
            return None
        if ("image" in lower_key or "base64" in lower_key) and len(value_str) > 500:
            return None
        if len(value_str) > 1800:
            return value_str[:900].rstrip() + " ...[сокращено]... " + value_str[-650:].lstrip()
        return value
    return value

def build_quality_review_prompt(
    structure: dict,
    user_prompt: str,
    doc_text: str,
    slide_count: int,
    style: str,
    tone: str,
) -> str:
    safe_structure = _strip_runtime_fields_for_llm(structure)
    draft_json = json.dumps(safe_structure, ensure_ascii=False, indent=2)
    doc_section = doc_text[:5000] if doc_text.strip() else "Документ не загружен. Проверяй фактические утверждения на осторожность формулировок; не придумывай новые цифры."
    return f"""Проверь и отредактируй презентацию как второй независимый LLM-контролёр.

ТЕМА ПОЛЬЗОВАТЕЛЯ:
{user_prompt}

ПАРАМЕТРЫ:
- Нужно ровно {slide_count} слайдов.
- Первый слайд: layout title.
- Последний слайд: layout conclusion.
- Тема оформления: {style_label(style)}.
- Тон: {tone}.
- Язык: русский.
- Не удаляй поля image_prompt, если они релевантны.
- Не добавляй imageData.
- Не меняй количество слайдов.
- Не добавляй неподтверждённые цифры или факты.
- Если источник не подтверждает факт, смягчи формулировку или убери число.
- Сделай текст компактным, чтобы он не съезжал в PPTX:
  title <= 78 символов;
  content bullets: до 5 пунктов, каждый до 110 символов;
  two_column: до 4 пунктов в каждой колонке;
  stats: до 3 карточек;
  conclusion <= 420 символов.

ИСХОДНЫЙ ДОКУМЕНТ/КОНТЕКСТ:
---
{doc_section}
---

ЧЕРНОВИК JSON:
{draft_json[:14000]}

Верни строго JSON такого вида:
{{
  "presentation_title": "...",
  "slides": [
    {{"layout":"title","title":"...","subtitle":"..."}}
  ],
  "quality_review": {{
    "score": 0,
    "summary": "...",
    "issues": [
      {{"slide": 1, "severity": "low|medium|high", "issue": "...", "fix": "..."}}
    ],
    "corrections": ["..."]
  }}
}}"""

def _merge_runtime_slide_fields(candidate: dict, original: dict) -> dict:
    """Возвращает runtime-поля после QA-LLM, чтобы не потерять картинки и image_prompt."""
    result = copy.deepcopy(candidate)
    old_slides = original.get("slides", []) or []
    new_slides = result.get("slides", []) or []
    for idx, slide in enumerate(new_slides):
        if idx >= len(old_slides):
            continue
        old = old_slides[idx]
        for key in ("imageData",):
            if old.get(key) and not slide.get(key):
                slide[key] = old[key]
        if old.get("image_prompt") and not slide.get("image_prompt"):
            slide["image_prompt"] = old.get("image_prompt")
    result["slides"] = new_slides
    return result

def _candidate_from_review_json(obj: dict, base: dict) -> dict:
    if not isinstance(obj, dict):
        raise ValueError("QA LLM вернула не JSON-объект")
    if isinstance(obj.get("slides"), list):
        return obj
    if isinstance(obj.get("corrected_presentation"), dict):
        return obj["corrected_presentation"]
    if isinstance(obj.get("presentation"), dict):
        return obj["presentation"]
    if isinstance(obj.get("corrected_slides"), list):
        candidate = copy.deepcopy(base)
        candidate["slides"] = obj["corrected_slides"]
        if isinstance(obj.get("quality_review"), dict):
            candidate["quality_review"] = obj["quality_review"]
        return candidate
    raise ValueError("QA LLM не вернула массив slides")

def review_and_refine_slide_structure(
    structure: dict,
    user_prompt: str,
    doc_text: str,
    slide_count: int,
    style: str,
    tone: str,
    rt_token: str = "",
) -> dict:
    """Второй проход: QA через LLM при наличии токена, иначе локальная чистка."""
    base = _normalize_slide_structure(copy.deepcopy(structure), slide_count)
    review_errors: list[str] = []

    if rt_token.strip():
        try:
            prompt = build_quality_review_prompt(base, user_prompt, doc_text, slide_count, style, tone)
            raw = _call_rt_llm(
                prompt,
                rt_token.strip(),
                system_prompt=QUALITY_SYSTEM_PROMPT,
                max_new_tokens=1536,
                temperature=0.08,
            )
            obj = _extract_json_object(raw)
            candidate = _candidate_from_review_json(obj, base)
            candidate = _normalize_slide_structure(candidate, slide_count)
            candidate = _merge_runtime_slide_fields(candidate, base)

            # Детерминированная проверка нужна даже после LLM, чтобы текст не переполнял PPTX.
            locally_checked = _local_quality_review_and_edit(candidate, user_prompt, doc_text, slide_count, style, tone)
            llm_report = obj.get("quality_review") if isinstance(obj.get("quality_review"), dict) else candidate.get("quality_review", {})
            local_report = locally_checked.get("quality_review", {})
            merged_issues = (llm_report.get("issues") or []) + (local_report.get("issues") or [])
            merged_corrections = (llm_report.get("corrections") or []) + (local_report.get("corrections") or [])
            score_values = [v for v in [llm_report.get("score"), local_report.get("score")] if isinstance(v, (int, float))]
            score = int(min(score_values)) if score_values else local_report.get("score", 90)

            locally_checked["quality_review"] = {
                "source": "RT LLM Quality Reviewer + локальный контроль",
                "score": score,
                "summary": llm_report.get("summary") or local_report.get("summary") or "Презентация проверена вторым LLM-агентом и локальными ограничителями.",
                "issues": merged_issues[:50],
                "corrections": merged_corrections[:50],
            }
            locally_checked["review_source"] = "RT LLM Quality Reviewer"
            return locally_checked
        except Exception as e:
            review_errors.append(f"RT QA LLM: {e}")

    reviewed = _local_quality_review_and_edit(base, user_prompt, doc_text, slide_count, style, tone)
    if review_errors:
        reviewed["review_warning"] = "Второй LLM-контролёр недоступен, применён локальный QA-контроль. " + " | ".join(review_errors)
    elif not rt_token.strip():
        reviewed["review_warning"] = "RT-токен не задан: вместо второго LLM-контролёра применён локальный QA-контроль."
    return reviewed

def build_single_slide_edit_prompt(
    structure: dict,
    slide_index: int,
    instruction: str,
    doc_text: str,
) -> str:
    slides = structure.get("slides", []) or []
    slide = _strip_runtime_fields_for_llm(slides[slide_index])
    context_slides = [
        {"index": i + 1, "layout": s.get("layout"), "title": s.get("title"), "plain_text": _slide_plain_text(s)[:500]}
        for i, s in enumerate(slides)
    ]
    return f"""Отредактируй только один слайд презентации.

Номер слайда: {slide_index + 1}
Инструкция пользователя:
{instruction or "Проверь слайд и исправь структуру, читаемость, достоверность и визуальный баланс."}

Исходный слайд:
{json.dumps(slide, ensure_ascii=False, indent=2)}

Контекст всей презентации:
{json.dumps(context_slides, ensure_ascii=False, indent=2)[:6500]}

Источник/документ:
---
{(doc_text or "Источник не загружен. Не добавляй новые факты и цифры без необходимости.")[:3500]}
---

Требования:
- Верни только исправленный JSON одного слайда в поле "slide".
- Не меняй остальные слайды.
- Сохрани layout, если инструкция явно не требует другого; для первого слайда layout title, для последнего conclusion.
- Текст должен помещаться: заголовок до 78 символов, пункты короткие, без перегруза.
- Не добавляй неподтвержденные цифры.
- Сохрани image_prompt, если он есть и подходит.
- Если пользователь просит фото/картинку/изображение, обязательно заполни image_prompt коротким английским промптом для 16:9.
- Не возвращай imageData: сервер сам сгенерирует файл изображения по image_prompt для выбранного слайда.

Формат:
{{
  "slide": {{"layout":"content","title":"...","bullets":["..."]}},
  "review": {{
    "issues": ["..."],
    "fixes": ["..."]
  }}
}}"""

def _slide_from_editor_payload(original: dict, payload: dict, slide_index: int, total: int) -> dict:
    slide = copy.deepcopy(original)
    if isinstance(payload.get("slide"), dict):
        slide.update(payload["slide"])

    if "layout" in payload and payload.get("layout"):
        slide["layout"] = _normalize_layout_name(payload.get("layout"))

    if "title" in payload:
        slide["title"] = _clean_text(payload.get("title"))

    body = payload.get("body")
    if body is not None:
        layout = _normalize_layout_name(slide.get("layout"))
        body_text = str(body or "").strip()
        if layout == "title":
            slide["subtitle"] = body_text
        elif layout == "quote":
            slide["quote"] = body_text
        elif layout in {"conclusion", "section_break", "section"}:
            slide["content"] = body_text
            slide["subtitle"] = body_text
        elif layout == "stats":
            slide["content"] = body_text
        elif layout == "two_column":
            lines = [x.strip(" •-\t") for x in body_text.splitlines() if x.strip(" •-\t")]
            half = max(1, (len(lines) + 1) // 2)
            slide["leftContent"] = lines[:half]
            slide["rightContent"] = lines[half:]
        else:
            slide["bullets"] = [x.strip(" •-\t") for x in body_text.splitlines() if x.strip(" •-\t")]

    if slide_index == 0:
        slide["layout"] = "title"
    if slide_index == total - 1:
        slide["layout"] = "conclusion"
    return slide

def edit_single_slide_with_reviewer(
    structure: dict,
    slide_index: int,
    instruction: str,
    rt_token: str,
    doc_text: str,
) -> tuple[dict, dict]:
    slides = structure.get("slides", []) or []
    if slide_index < 0 or slide_index >= len(slides):
        raise ValueError("Неверный номер слайда")

    report: dict[str, Any] = {
        "source": "Локальный редактор",
        "issues": [],
        "fixes": [],
    }

    if rt_token.strip():
        try:
            prompt = build_single_slide_edit_prompt(structure, slide_index, instruction, doc_text)
            raw = _call_rt_llm(
                prompt,
                rt_token.strip(),
                system_prompt=SINGLE_SLIDE_EDITOR_SYSTEM_PROMPT,
                max_new_tokens=1200,
                temperature=0.12,
            )
            obj = _extract_json_object(raw)
            candidate = obj.get("slide") if isinstance(obj.get("slide"), dict) else obj
            if not isinstance(candidate, dict):
                raise ValueError("Редактор слайда не вернул объект slide")
            original = slides[slide_index]
            if original.get("imageData") and not candidate.get("imageData"):
                candidate["imageData"] = original["imageData"]
            if original.get("image_prompt") and not candidate.get("image_prompt"):
                candidate["image_prompt"] = original["image_prompt"]
            cleaned, issues, fixes = _sanitize_slide_for_quality(
                candidate,
                slide_index,
                len(slides),
                _doc_keyword_set(doc_text),
                doc_text,
            )
            if isinstance(obj.get("review"), dict):
                report.update(obj["review"])
            report["source"] = "RT LLM редактор одного слайда"
            report["issues"] = (report.get("issues") or []) + issues
            report["fixes"] = (report.get("fixes") or []) + fixes
            slides[slide_index] = cleaned
            structure["slides"] = slides
            return structure, report
        except Exception as e:
            report["issues"].append(f"LLM-редактор не сработал: {e}. Применены только локальные правки.")

    cleaned, issues, fixes = _sanitize_slide_for_quality(
        slides[slide_index],
        slide_index,
        len(slides),
        _doc_keyword_set(doc_text),
        doc_text,
    )
    slides[slide_index] = cleaned
    structure["slides"] = slides
    report["issues"].extend(issues)
    report["fixes"].extend(fixes)
    return structure, report

def _session_json_path(session_id: str) -> Path:
    if not re.fullmatch(r"[0-9a-fA-F-]{20,80}", session_id or ""):
        raise HTTPException(400, "Некорректный session_id")
    return WORK_DIR / f"{session_id}.json"

def _session_pptx_path(session_id: str) -> Path:
    if not re.fullmatch(r"[0-9a-fA-F-]{20,80}", session_id or ""):
        raise HTTPException(400, "Некорректный session_id")
    return WORK_DIR / f"{session_id}.pptx"

def load_session_structure(session_id: str) -> dict:
    json_path = _session_json_path(session_id)
    if not json_path.exists():
        raise HTTPException(404, "Сессия презентации не найдена")
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)

def save_session_structure(session_id: str, structure: dict, style: str) -> None:
    output_path = str(_session_pptx_path(session_id))
    if not build_pptx(structure, output_path, style):
        raise HTTPException(500, "Ошибка пересборки PPTX")
    with open(_session_json_path(session_id), "w", encoding="utf-8") as f:
        json.dump(structure, f, ensure_ascii=False, indent=2)

def _session_context(structure: dict, session_id: str = "") -> dict:
    """Объединяет сохранённый контекст и runtime-секреты текущей сессии."""
    context = structure.get("_presentai_context") if isinstance(structure.get("_presentai_context"), dict) else {}
    runtime = SESSION_SECRETS.get(session_id, {}) if session_id else {}
    merged = {**context, **runtime}
    if "style" in merged:
        merged["style"] = normalize_style(merged["style"])
    return merged

def generate_slide_structure(user_prompt: str, doc_text: str,
                               slide_count: int, style: str, tone: str,
                               rt_token: str = "") -> dict:
    """Пробует RT LLM, затем Anthropic fallback, затем локальную демо-структуру."""
    prompt = build_generation_prompt(user_prompt, doc_text, slide_count, style, tone)
    errors = []

    if rt_token.strip():
        try:
            rt_text = _call_rt_llm(prompt, rt_token.strip())
            try:
                structure = _normalize_slide_structure(_extract_json_object(rt_text), slide_count)
            except Exception:
                # Иногда LLM отвечает осмысленным текстом вместо JSON; сохраняем демо-работоспособность.
                structure = _normalize_slide_structure(_structure_from_plain_text(rt_text, user_prompt, slide_count, style, tone), slide_count)
                structure["generation_warning"] = "RT LLM вернула текст не в JSON-формате, слайды собраны из текста ответа."
            structure["generation_source"] = "RT LLM"
            return structure
        except Exception as e:
            errors.append(f"RT LLM: {e}")

    if os.getenv("USE_ANTHROPIC_FALLBACK") == "1" and os.getenv("ANTHROPIC_API_KEY"):
        try:
            structure = _normalize_slide_structure(_extract_json_object(_call_anthropic(prompt)), slide_count)
            structure["generation_source"] = "Anthropic"
            return structure
        except Exception as e:
            errors.append(f"Anthropic: {e}")

    structure = _fallback_slide_structure(user_prompt, doc_text, slide_count, style, tone)
    structure["generation_source"] = "Локальный fallback"
    if errors:
        structure["generation_warning"] = "Использован локальный fallback. " + " | ".join(errors)
    else:
        structure["generation_warning"] = "Использован локальный fallback: API-токен LLM не задан."
    return _normalize_slide_structure(structure, slide_count)

# ── Генерация изображений через RT API ────────────────────────────────────────

LAST_IMAGE_ERROR = ""


def _set_image_error(message: object) -> None:
    """Запоминает последнюю ошибку image API, чтобы UI показал понятную причину."""
    global LAST_IMAGE_ERROR
    text = _clean_text(message)
    LAST_IMAGE_ERROR = text[:900]


def _normalize_rt_image_service(service: object) -> str:
    """Нормализует имя сервиса из UI к идентификаторам RT endpoints."""
    value = str(service or "").strip().lower().replace("-", "_").replace(" ", "")
    if value in {"ya", "yandex", "yandexart", "yandex_art", "yaart", "арт", "yandexарт"}:
        return "yaArt"
    return "sd"


def _rt_image_service_candidates(service: object) -> list[str]:
    """Подбирает варианты serviceType: RT download чувствителен к точному значению."""
    normalized = _normalize_rt_image_service(service)
    candidates: list[str] = []

    def add(value: object) -> None:
        s = str(value or "").strip()
        if s and s not in candidates:
            candidates.append(s)

    add(normalized)
    if normalized == "yaArt":
        add("yaArt")
        # В описании Yandex ART встречается неоднозначный serviceType, поэтому
        # для /download держим запасной SD-вариант, хотя job создаём через /ya/image.
        add("sd")
    else:
        add("sd")
        add("yaArt")
    return candidates


def _guess_image_mime(data: bytes) -> str:
    """Определяет MIME изображения по сигнатуре байтов без внешних зависимостей."""
    data = data or b""
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"GIF87a") or data.startswith(b"GIF89a"):
        return "image/gif"
    if data.startswith(b"RIFF") and b"WEBP" in data[:16]:
        return "image/webp"
    if data.startswith(b"BM"):
        return "image/bmp"
    if data[:4] in (b"II*\x00", b"MM\x00*"):
        return "image/tiff"
    return "image/png"


def _is_supported_image_bytes(data: bytes) -> bool:
    data = data or b""
    return (
        data.startswith(b"\x89PNG\r\n\x1a\n")
        or data.startswith(b"\xff\xd8\xff")
        or data.startswith(b"GIF87a")
        or data.startswith(b"GIF89a")
        or (data.startswith(b"RIFF") and b"WEBP" in data[:16])
        or data.startswith(b"BM")
        or data[:4] in (b"II*\x00", b"MM\x00*")
    )


def _normalize_generated_image_bytes(data: bytes) -> bytes:
    """Возвращает байты, которые python-pptx сможет встроить; экзотику переводит в PNG."""
    if not data:
        return b""
    if data.startswith(b"\x89PNG\r\n\x1a\n") or data.startswith(b"\xff\xd8\xff") or data.startswith(b"GIF87a") or data.startswith(b"GIF89a") or data.startswith(b"BM"):
        return data
    if _is_supported_image_bytes(data):
        try:
            from PIL import Image
            out = io.BytesIO()
            with Image.open(io.BytesIO(data)) as im:
                if im.mode not in ("RGB", "RGBA"):
                    im = im.convert("RGBA")
                im.save(out, format="PNG")
            return out.getvalue()
        except Exception:
            return data
    return b""


def _image_bytes_from_text(text: object) -> bytes:
    """Достаёт изображение из data URL или сырой base64-строки."""
    if not isinstance(text, str):
        return b""
    value = text.strip()
    if not value:
        return b""
    if value.startswith("data:image/") or value.startswith("image/"):
        raw = value.split(",", 1)[1] if "," in value else ""
    elif len(value) > 120 and re.fullmatch(r"[A-Za-z0-9+/=\s]+", value):
        raw = value
    else:
        return b""
    raw = re.sub(r"\s+", "", raw)
    try:
        decoded = base64.b64decode(raw, validate=False)
    except Exception:
        return b""
    return _normalize_generated_image_bytes(decoded)


def _image_bytes_from_json_payload(obj: object) -> bytes:
    """RT иногда заворачивает файл/base64 в JSON; рекурсивно ищем данные картинки."""
    if isinstance(obj, str):
        return _image_bytes_from_text(obj)
    if isinstance(obj, list):
        for item in obj:
            found = _image_bytes_from_json_payload(item)
            if found:
                return found
    if isinstance(obj, dict):
        preferred_keys = (
            "image", "imageData", "image_data", "base64", "data", "content",
            "file", "fileData", "bytes", "payload", "result",
        )
        for key in preferred_keys:
            if key in obj:
                found = _image_bytes_from_json_payload(obj.get(key))
                if found:
                    return found
        for value in obj.values():
            found = _image_bytes_from_json_payload(value)
            if found:
                return found
    return b""


def _urls_from_json_payload(obj: object) -> list[str]:
    """Собирает URL-поля из JSON-ответа /download."""
    urls: list[str] = []

    def add(value: object) -> None:
        if not isinstance(value, str):
            return
        s = value.strip()
        if not s:
            return
        if s.startswith("http://") or s.startswith("https://") or s.startswith("/api/") or s.startswith("/download") or s.startswith("download"):
            if s not in urls:
                urls.append(s)

    if isinstance(obj, str):
        add(obj)
    elif isinstance(obj, list):
        for item in obj:
            urls.extend(x for x in _urls_from_json_payload(item) if x not in urls)
    elif isinstance(obj, dict):
        for key in ("url", "href", "link", "download", "downloadUrl", "fileUrl", "path"):
            add(obj.get(key))
        for value in obj.values():
            urls.extend(x for x in _urls_from_json_payload(value) if x not in urls)
    return urls


def _decode_json_bytes(data: bytes) -> object:
    try:
        return json.loads((data or b"").decode("utf-8", errors="replace"))
    except Exception:
        return None


def _extract_rt_message(data: object) -> dict:
    """Находит вложенный объект message в ответе RT image API."""
    if isinstance(data, dict):
        message = data.get("message")
        if isinstance(message, dict):
            return message
        for value in data.values():
            found = _extract_rt_message(value)
            if found:
                return found
    if isinstance(data, list):
        for item in data:
            found = _extract_rt_message(item)
            if found:
                return found
    return {}


def _short_rt_body(resp: object, limit: int = 280) -> str:
    try:
        text = _decode_rt_response_body(resp, limit=limit)
    except Exception:
        text = str(getattr(resp, "text", "") or "")[:limit]
    return _clean_text(text)


def _download_url_image(url: str, token: str, req_module: object) -> bytes:
    """Скачивает изображение по URL из RT JSON, сохраняя Authorization."""
    if not url:
        return b""
    from urllib.parse import urljoin
    if url.startswith("/"):
        full_url = "https://ai.rt.ru" + url
    elif url.startswith("download"):
        full_url = urljoin(RT_API_BASE.rstrip("/") + "/", url)
    else:
        full_url = url
    headers = {"Authorization": f"Bearer {token}", "Accept": "image/png,image/*,*/*"}
    try:
        resp = req_module.get(full_url, headers=headers, timeout=RT_IMAGE_REQUEST_TIMEOUT)
        if resp.status_code >= 400:
            _set_image_error(f"GET URL {resp.status_code}: {_short_rt_body(resp)}")
            return b""
        data = resp.content or b""
        image = _normalize_generated_image_bytes(data)
        if image:
            return image
        parsed = _decode_json_bytes(data)
        if parsed is not None:
            return _image_bytes_from_json_payload(parsed)
    except Exception as e:
        _set_image_error(f"GET URL error: {e}")
    return b""


def _image_from_download_response(resp: object, token: str, req_module: object) -> bytes:
    """Разбирает /download как бинарную картинку, JSON с base64 или URL."""
    content = getattr(resp, "content", b"") or b""
    image = _normalize_generated_image_bytes(content)
    if image:
        return image

    parsed = _decode_json_bytes(content)
    if parsed is not None:
        image = _image_bytes_from_json_payload(parsed)
        if image:
            return image
        for url in _urls_from_json_payload(parsed)[:4]:
            image = _download_url_image(url, token, req_module)
            if image:
                return image
    return b""


def _download_rt_image_file(
    msg_id: object,
    token: str,
    service_type: object,
    preferred_service: object,
    req_module: object,
    timeout_seconds: Optional[float] = None,
) -> tuple[bytes, str]:
    """Опрашивает RT /download, пока реальный файл изображения не станет доступен."""
    services: list[str] = []
    for svc in [service_type, preferred_service, *_rt_image_service_candidates(service_type), *_rt_image_service_candidates(preferred_service)]:
        normalized = _normalize_rt_image_service(svc)
        for candidate in [str(svc or "").strip(), normalized]:
            if candidate and candidate not in services:
                services.append(candidate)
    if not services:
        services = ["sd", "yaArt"]

    image_types = ["png", "jpeg"]
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
        "Accept": "image/png,image/jpeg,image/*;q=0.9,application/json;q=0.6,*/*;q=0.2",
    }
    deadline = time.time() + max(1, min(RT_IMAGE_DOWNLOAD_TIMEOUT, timeout_seconds or RT_IMAGE_DOWNLOAD_TIMEOUT))
    attempt = 0
    last_error = ""

    while time.time() < deadline:
        for svc in services:
            for image_type in image_types:
                remaining = deadline - time.time()
                if remaining <= 0:
                    break
                try:
                    resp = req_module.get(
                        f"{RT_API_BASE}/download",
                        params={"id": msg_id, "serviceType": svc, "imageType": image_type},
                        headers=headers,
                        timeout=max(1, min(RT_IMAGE_DOWNLOAD_REQUEST_TIMEOUT, remaining)),
                    )
                    if resp.status_code == 429:
                        last_error = f"download id={msg_id}, serviceType={svc}: HTTP 429 (rate limit)"
                        remaining_after = deadline - time.time()
                        if remaining_after > RT_IMAGE_RATELIMIT_PAUSE:
                            time.sleep(RT_IMAGE_RATELIMIT_PAUSE)
                        continue
                    if resp.status_code >= 400:
                        last_error = f"download id={msg_id}, serviceType={svc}, imageType={image_type}: HTTP {resp.status_code} {_short_rt_body(resp)}"
                        continue
                    image = _image_from_download_response(resp, token, req_module)
                    if image:
                        return image, ""
                    content_type = str(resp.headers.get("Content-Type", "")) if hasattr(resp, "headers") else ""
                    body = _short_rt_body(resp, 220)
                    last_error = f"download id={msg_id}, serviceType={svc}: ответ не является изображением; Content-Type={content_type}; body={body}"
                except Exception as e:
                    last_error = f"download id={msg_id}, serviceType={svc}: {e}"
            if time.time() >= deadline:
                break
        attempt += 1
        if time.time() < deadline:
            time.sleep(min(1.2 + attempt * 0.35, 3.0, max(0.1, deadline - time.time())))

    return b"", last_error or f"download id={msg_id}: файл не появился за время ожидания"


def _post_rt_image_job(prompt: str, token: str, service: str, translate: bool, req_module: object) -> tuple[object, str, bytes, str]:
    """Создаёт RT image job и возвращает message id, serviceType, прямую картинку или ошибку."""
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json; charset=utf-8",
        "Accept": "application/json,*/*",
    }
    request_id = str(uuid.uuid4())
    seed = int(time.time() * 1000) % 2147483647

    if service == "yaArt":
        payload = {
            "uuid": request_id,
            "image": {
                "request": prompt,
                "seed": seed,
                "translate": translate,
                "model": "yandex-art",
                "aspect": "16:9",
            },
        }
        endpoint = f"{RT_API_BASE}/ya/image"
    else:
        payload = {
            "uuid": request_id,
            "sdImage": {
                "request": prompt,
                "seed": seed,
                "translate": translate,
            },
        }
        endpoint = f"{RT_API_BASE}/sd/img"

    last_err = ""
    for attempt in range(max(1, RT_IMAGE_POST_RETRIES + 1)):
        try:
            resp = req_module.post(endpoint, json=payload, headers=headers, timeout=RT_IMAGE_POST_TIMEOUT)
            if resp.status_code == 429:
                last_err = f"POST {endpoint} HTTP 429 (rate limit): {_short_rt_body(resp)}"
                if attempt < RT_IMAGE_POST_RETRIES:
                    time.sleep(RT_IMAGE_RATELIMIT_PAUSE)
                    continue
                return None, service, b"", last_err
            if resp.status_code >= 400:
                return None, service, b"", f"POST {endpoint} HTTP {resp.status_code}: {_short_rt_body(resp)}"
            content = resp.content or b""
            direct_image = _normalize_generated_image_bytes(content)
            if direct_image:
                return None, service, direct_image, ""
            try:
                data = resp.json()
            except Exception:
                data = _decode_json_bytes(content)
            if data is None:
                return None, service, b"", f"POST {endpoint}: ответ не JSON и не изображение: {_short_rt_body(resp)}"
            direct_image = _image_bytes_from_json_payload(data)
            if direct_image:
                return None, service, direct_image, ""
            message = _extract_rt_message(data)
            msg_id = message.get("id") or message.get("messageId") or message.get("message_id")
            returned_service = message.get("serviceType") or service
            if not msg_id:
                return None, returned_service, b"", f"POST {endpoint}: в ответе нет message.id; фрагмент={str(data)[:320]}"
            return msg_id, returned_service, b"", ""
        except Exception as e:
            last_err = f"POST {endpoint}: {e}"
            # На таймаут сразу возвращаем ошибку — ждать смысла нет.
            if "timed out" in str(e).lower() or "timeout" in str(e).lower():
                return None, service, b"", last_err
    return None, service, b"", last_err


def _prepare_image_prompt(prompt: object) -> str:
    """Сжимает image prompt до компактного и безопасного для RT endpoints вида."""
    text = _clean_text(prompt)
    if not text:
        text = "professional realistic 16:9 editorial presentation image, high quality, no text overlay"
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text).strip()
    return _clip_text(text, 900)


def generate_image_rt(prompt: str, token: str, service: str = "sd") -> Optional[bytes]:
    """Генерирует изображение через RT API (Stable Diffusion или Yandex ART).

    RT image API сначала возвращает message id, а бинарный файл нужно забрать
    через /download. На практике /download может временно вернуть JSON, пустое
    тело или другой serviceType, поэтому здесь есть несколько fallback-попыток.
    """
    _set_image_error("")
    token = (token or "").strip()
    if not token:
        _set_image_error("RT-токен не задан")
        return None

    try:
        import requests as req
    except Exception as e:
        _set_image_error(f"Не удалось импортировать requests: {e}")
        return None

    image_prompt = _prepare_image_prompt(prompt)
    deadline = time.time() + RT_IMAGE_SINGLE_TIMEOUT
    requested_service = _normalize_rt_image_service(service)
    service_order = [requested_service]
    other_service = "sd" if requested_service == "yaArt" else "yaArt"
    if RT_IMAGE_ALLOW_BACKEND_FALLBACK and other_service not in service_order:
        service_order.append(other_service)

    errors: list[str] = []
    for svc in service_order:
        # translate=true нужен только для русскоязычного image prompt; для английских
        # prompt-ов лишняя попытка обычно просто удваивает время ожидания.
        translate_options = (False, True) if re.search(r"[А-Яа-яЁё]", image_prompt) else (False,)
        for translate in translate_options:
            if time.time() >= deadline:
                errors.append(f"{svc}: превышен лимит {RT_IMAGE_SINGLE_TIMEOUT} секунд на одно изображение")
                break
            msg_id, returned_service, direct_image, post_error = _post_rt_image_job(
                image_prompt, token, svc, translate, req
            )
            if direct_image:
                return direct_image
            if post_error:
                errors.append(post_error)
                continue
            if not msg_id:
                errors.append(f"{svc}: RT API не вернул id задания изображения")
                continue

            remaining = max(1, deadline - time.time())
            image, download_error = _download_rt_image_file(
                msg_id,
                token,
                returned_service,
                svc,
                req,
                timeout_seconds=remaining,
            )
            if image:
                return image
            errors.append(download_error or f"{svc}: /download не вернул файл изображения")
            if time.time() >= deadline:
                break

    _set_image_error(" | ".join(x for x in errors if x) or "RT API изображений не вернул файл")
    print(f"Image generation failed: {LAST_IMAGE_ERROR}")
    return None


def image_bytes_to_base64(data: bytes) -> str:
    """Возвращает data URL, безопасный для браузера и PPTX."""
    data = _normalize_generated_image_bytes(data)
    if not data:
        return ""
    mime = _guess_image_mime(data)
    return f"data:{mime};base64," + base64.b64encode(data).decode()


def _normalize_image_data_url(value: object) -> str:
    """Принимает старый `image/png;base64,...` и нормальный `data:image/...` формат."""
    if not isinstance(value, str):
        return ""
    text = value.strip()
    if not text:
        return ""
    if text.startswith("data:image/"):
        return text
    if text.startswith("image/") and ";base64," in text:
        return "data:" + text
    if text.startswith("image/") and "," in text:
        return "data:" + text
    if text.startswith("/") or text.startswith("http://") or text.startswith("https://"):
        return text
    return "data:image/png;base64," + text


def _decode_image_data(value: object) -> bytes:
    """Безопасно декодирует imageData; при ошибке возвращает пустые байты."""
    text = _normalize_image_data_url(value)
    if not text or text.startswith("http") or text.startswith("/"):
        return b""
    raw = text.split(",", 1)[1] if "," in text else text
    raw = re.sub(r"\s+", "", raw)
    try:
        decoded = base64.b64decode(raw, validate=False)
    except Exception:
        return b""
    return _normalize_generated_image_bytes(decoded)


_IMAGE_REQUEST_RE = re.compile(
    r"(фот|фото|картин|изображ|иллюстр|визуал|логотип|logo|photo|image|picture|illustration|visual)",
    re.IGNORECASE,
)


def _looks_like_image_request(text: object) -> bool:
    return bool(_IMAGE_REQUEST_RE.search(str(text or "")))


def _fallback_image_prompt_for_slide(slide: dict, topic: str = "", instruction: str = "") -> str:
    """Создаёт image_prompt, если LLM не вернула его для слайда."""
    title = _clean_text(slide.get("title") or topic or "presentation slide")
    plain = _slide_plain_text(slide)
    source = " ".join(x for x in [instruction, title, plain[:220], topic] if _clean_text(x))
    source = _clip_text(source, 420)
    return (
        "professional realistic 16:9 editorial presentation image, "
        f"topic: {source}, high quality, clean composition, no text overlay"
    )


def _ensure_slide_image_prompt(slide: dict, topic: str = "", instruction: str = "", force: bool = False) -> bool:
    """Гарантирует наличие image_prompt, когда пользователь запросил генерацию картинок."""
    if not isinstance(slide, dict):
        return False
    if _clean_text(slide.get("image_prompt")):
        return False
    layout = _normalize_layout_name(slide.get("layout"))
    if force or layout not in {"title", "conclusion", "section_break", "section"}:
        slide["image_prompt"] = _fallback_image_prompt_for_slide(slide, topic, instruction)
        return True
    return False


def _generate_slide_image_if_needed(
    slide: dict,
    token: str,
    service: str,
    topic: str = "",
    instruction: str = "",
    force: bool = False,
) -> tuple[bool, str]:
    """Генерирует и прикрепляет imageData только для одного нужного слайда."""
    if not isinstance(slide, dict):
        return False, "Слайд не является объектом JSON."
    token = (token or "").strip()
    if not token:
        return False, "RT-токен не задан, изображение не сгенерировано."

    image_requested = force or _looks_like_image_request(instruction) or bool(_clean_text(slide.get("image_prompt")))
    if not image_requested:
        return False, "Изображение для этого слайда не запрошено."

    _ensure_slide_image_prompt(slide, topic, instruction, force=True)
    if slide.get("imageData") and not force:
        slide["imageData"] = _normalize_image_data_url(slide.get("imageData"))
        return False, "На слайде уже есть изображение."

    img_bytes = generate_image_rt(_clean_text(slide.get("image_prompt")), token, service or "sd")
    if not img_bytes:
        detail = _clean_text(LAST_IMAGE_ERROR)
        if detail:
            return False, f"RT API изображений не вернул файл: {detail}"
        return False, "RT API изображений не вернул файл."
    if not _is_supported_image_bytes(img_bytes):
        preview = img_bytes[:120].decode("utf-8", errors="replace")
        return False, f"RT API вернул не изображение: {preview}"

    slide["imageData"] = image_bytes_to_base64(img_bytes)
    if not slide["imageData"]:
        return False, "Изображение получено, но не удалось подготовить его для PPTX."
    return True, "Изображение сгенерировано и прикреплено к выбранному слайду."

# ── Генерация PPTX ────────────────────────────────────────────────────────────

THEMES = {
    DEFAULT_STYLE: {
        "bg": "0A0F1E",
        "bg2": "1A0A3E",
        "title": "FFFFFF",
        "text": "D0D5DD",
        "accent": "7700FF",
        "accent2": "C026D3",
    },
    "modern": {
        "bg": "0F172A",
        "bg2": "1E3A5F",
        "title": "F8FAFC",
        "text": "CBD5E1",
        "accent": "38BDF8",
        "accent2": "818CF8",
    },
    "corporate": {
        "bg": "F1F5F9",
        "bg2": "DBEAFE",
        "title": "0F172A",
        "text": "334155",
        "accent": "2563EB",
        "accent2": "0EA5E9",
    },
    "minimal": {
        "bg": "FFFFFF",
        "bg2": "F1F5F9",
        "title": "111827",
        "text": "4B5563",
        "accent": "6366F1",
        "accent2": "8B5CF6",
    },
    "tech": {
        "bg": "071A0E",
        "bg2": "0D2B18",
        "title": "ECFDF5",
        "text": "BBF7D0",
        "accent": "22C55E",
        "accent2": "06B6D4",
    },
    "creative": {
        "bg": "1E0A35",
        "bg2": "3D0F52",
        "title": "FFF7ED",
        "text": "F5D0FE",
        "accent": "F97316",
        "accent2": "EC4899",
    },
}

def _xml_text(value: object) -> str:
    return escape(str(value or ""))

def _slide_lines(slide: dict) -> list[str]:
    layout = slide.get("layout", "content")
    lines: list[str] = []

    if layout == "title":
        lines.append(slide.get("subtitle", ""))
    elif layout == "two_column":
        lines.append(slide.get("leftTitle", ""))
        lines.extend(f"- {item}" for item in slide.get("leftContent", []))
        lines.append("")
        lines.append(slide.get("rightTitle", ""))
        lines.extend(f"- {item}" for item in slide.get("rightContent", []))
    elif layout == "stats":
        for stat in slide.get("stats", []):
            lines.append(f"{stat.get('value', '')} - {stat.get('label', '')}")
        if slide.get("content"):
            lines.append("")
            lines.append(slide.get("content", ""))
    elif layout == "quote":
        lines.append(slide.get("quote", ""))
        if slide.get("title"):
            lines.append(slide.get("title", ""))
    elif layout in {"section_break", "conclusion"}:
        lines.append(slide.get("subtitle", slide.get("content", "")))
    else:
        lines.extend(f"- {item}" for item in slide.get("bullets", []))
        if slide.get("content"):
            lines.append(slide.get("content", ""))

    return [line for line in lines if line is not None]

def _text_shape(shape_id: int, x: int, y: int, w: int, h: int,
                lines: list[str], font_size: int, color: str, bold: bool = False) -> str:
    paragraphs = []
    for line in lines or [""]:
        paragraphs.append(f"""
        <a:p>
          <a:r>
            <a:rPr lang="ru-RU" sz="{font_size}" b="{1 if bold else 0}">
              <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
            </a:rPr>
            <a:t>{_xml_text(line)}</a:t>
          </a:r>
        </a:p>""")

    return f"""
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="{shape_id}" name="TextBox {shape_id}"/>
        <p:cNvSpPr txBox="1"/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{x}" y="{y}"/>
          <a:ext cx="{w}" cy="{h}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="top"/>
        <a:lstStyle/>
        {''.join(paragraphs)}
      </p:txBody>
    </p:sp>"""

def _slide_xml(slide: dict, index: int, colors: dict) -> str:
    title = slide.get("title") or slide.get("presentation_title") or f"Слайд {index}"
    body_lines = _slide_lines(slide)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:bg>
      <p:bgPr><a:solidFill><a:srgbClr val="{colors['bg']}"/></a:solidFill></p:bgPr>
    </p:bg>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      {_text_shape(2, 650000, 650000, 10850000, 1300000, [title], 3600, colors["title"], True)}
      {_text_shape(3, 700000, 2150000, 10650000, 3900000, body_lines, 2100, colors["text"])}
      {_text_shape(4, 700000, 6350000, 2500000, 450000, [f"{index}"], 1400, colors["accent"], True)}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""

def build_pptx(slide_data: dict, output_path: str, theme: str) -> bool:
    """Собирает PPTX с градиентами, акцентными фигурами и изображениями."""
    slides_list = slide_data.get("slides") or []
    if not slides_list:
        return False

    colors = THEMES.get(normalize_style(theme), THEMES[DEFAULT_STYLE])

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        bg_c  = colors["bg"]
        bg2_c = colors.get("bg2", bg_c)
        ti_c  = colors["title"]
        tx_c  = colors["text"]
        ac_c  = colors["accent"]
        ac2_c = colors.get("accent2", ac_c)
        NS_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"

        def rgb(h: str) -> RGBColor:
            return RGBColor.from_string(h)

        def _grad_bg(slide, c1: str, c2: str) -> None:
            """Ставит градиентный фон через XML, потому что python-pptx не даёт удобного API."""
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(c1)
            bgPr = slide._element.find('.//' + qn('p:bgPr'))
            if bgPr is None:
                return
            for ch in list(bgPr):
                bgPr.remove(ch)
            bgPr.append(etree.fromstring(
                f'<a:gradFill xmlns:a="{NS_A}" rotWithShape="1">'
                f'<a:gsLst>'
                f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
                f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
                f'</a:gsLst>'
                f'<a:lin ang="16200000" scaled="0"/>'
                f'</a:gradFill>'
            ))

        def _shape(sl, sid: int, x, y, w, h, clr: str, alpha: int = 100):
            """Добавляет фигуру с прозрачностью; alpha 100 означает полностью непрозрачно."""
            sh = sl.shapes.add_shape(sid, Inches(x), Inches(y), Inches(w), Inches(h))
            sh.line.fill.background()
            sh.fill.solid()
            sh.fill.fore_color.rgb = rgb(clr)
            if alpha < 100:
                sp = sh._element.find(qn('p:spPr'))
                sf = sp.find('.//' + qn('a:solidFill'))
                if sf is not None:
                    sc = sf.find(qn('a:srgbClr'))
                    if sc is not None:
                        ae = etree.SubElement(sc, qn('a:alpha'))
                        ae.set('val', str(alpha * 1000))
            return sh

        def rect(sl, x, y, w, h, clr, a=100):
            return _shape(sl, 1, x, y, w, h, clr, a)

        def oval(sl, x, y, w, h, clr, a=15):
            return _shape(sl, 9, x, y, w, h, clr, a)

        def txt(sl, x, y, w, h, text, sz=22, bold=False, clr=None,
                align='left', wrap=True, italic=False):
            amap = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
            box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = wrap
            lines = text.split('\n') if isinstance(text, str) else list(text or [''])
            for i, ln in enumerate(lines or ['']):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(ln)
                p.alignment = amap.get(align, PP_ALIGN.LEFT)
                p.font.size = Pt(sz)
                p.font.bold = bold
                p.font.italic = italic
                p.font.name = 'Arial'
                p.font.color.rgb = rgb(clr or tx_c)
            return box

        def fit_sz(text, base=22, min_sz=13, capacity=360, line_capacity=8):
            raw = "\n".join(str(x) for x in text) if isinstance(text, (list, tuple)) else str(text or "")
            chars = len(raw)
            lines_n = max(1, raw.count("\n") + 1)
            overflow = max(0, chars - capacity)
            penalty = (overflow // 90) + max(0, lines_n - line_capacity)
            return max(min_sz, int(base - penalty))

        def decor(sl):
            """Добавляет полупрозрачные декоративные круги по углам."""
            oval(sl,  9.5, -1.5, 6.5, 6.5, ac_c,  a=10)
            oval(sl, -2.0,  5.0, 5.5, 5.5, ac2_c, a=7)

        def slide_num(sl, n):
            txt(sl, 12.1, 7.05, 1.1, 0.3, str(n), sz=11, bold=True, clr=ac_c, align='right')

        def add_slide_image(sl, item, x, y, w, h, panel=True):
            """Встраивает imageData в PPTX с сохранением пропорций."""
            blob = _decode_image_data(item.get('imageData'))
            if not blob:
                return False
            if not _is_supported_image_bytes(blob):
                print('Image embed skipped: imageData is not a supported image file')
                return False
            if panel:
                rect(sl, x - 0.03, y - 0.03, w + 0.06, h + 0.06, 'FFFFFF', a=8)
            px, py, pw, ph = x, y, w, h
            try:
                from PIL import Image
                with Image.open(io.BytesIO(blob)) as im:
                    iw, ih = im.size
                if iw > 0 and ih > 0:
                    img_ratio = iw / ih
                    box_ratio = w / h
                    if img_ratio >= box_ratio:
                        pw = w
                        ph = w / img_ratio
                        py = y + (h - ph) / 2
                    else:
                        ph = h
                        pw = h * img_ratio
                        px = x + (w - pw) / 2
            except Exception as image_size_error:
                print(f'Image size detection failed: {image_size_error}')
            try:
                sl.shapes.add_picture(io.BytesIO(blob), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))
                return True
            except Exception as embed_error:
                print(f'Image embed failed: {embed_error}')
                return False

        for idx, item in enumerate(slides_list, start=1):
            sl = prs.slides.add_slide(blank)
            layout = item.get('layout', 'content')
            title  = item.get('title') or f'Слайд {idx}'

            _grad_bg(sl, bg_c, bg2_c)
            decor(sl)
            image_placed = False

            if layout == 'title':
                rect(sl, 0, 6.7, 13.333, 0.8, ac_c,  a=85)
                rect(sl, 0, 6.7, 13.333, 0.8, ac2_c, a=25)
                rect(sl, 0, 0, 0.3, 6.7, ac_c, a=90)
                rect(sl, 0.55, 1.6, 12.3, 3.9, ac_c, a=6)
                txt(sl, 0.9, 1.85, 11.5, 2.1, title, sz=fit_sz(title, 46, 34, 95, 2), bold=True, clr=ti_c)
                rect(sl, 0.9, 4.0, 3.8, 0.07, ac_c)
                sub = item.get('subtitle', '')
                if sub:
                    txt(sl, 0.9, 4.25, 11.0, 1.3, sub, sz=fit_sz(sub, 23, 17, 165, 3), clr=tx_c)

            elif layout == 'two_column':
                rect(sl, 0, 0, 13.333, 0.13, ac_c)
                txt(sl, 0.5, 0.28, 12.3, 0.85, title, sz=fit_sz(title, 28, 22, 90, 2), bold=True, clr=ti_c)
                rect(sl, 0.45, 1.4, 5.9, 5.75, ac_c,  a=10)
                rect(sl, 0.45, 1.4, 5.9, 0.14, ac_c,  a=90)
                txt(sl, 0.65, 1.65, 5.5, 0.55, item.get('leftTitle', ''), sz=18, bold=True, clr=ac_c)
                left_text = '\n'.join(f'• {x}' for x in item.get('leftContent', []))
                txt(sl, 0.65, 2.3,  5.5, 4.6,
                    left_text, sz=fit_sz(left_text, 17, 13, 410, 7), clr=tx_c)
                rect(sl, 6.5,  1.4, 0.07, 5.75, ac_c,  a=22)
                rect(sl, 6.72, 1.4, 6.1,  5.75, ac2_c, a=10)
                rect(sl, 6.72, 1.4, 6.1,  0.14, ac2_c, a=90)
                txt(sl, 6.92, 1.65, 5.7, 0.55, item.get('rightTitle', ''), sz=18, bold=True, clr=ac2_c)
                right_text = '\n'.join(f'• {x}' for x in item.get('rightContent', []))
                txt(sl, 6.92, 2.3,  5.7, 4.6,
                    right_text, sz=fit_sz(right_text, 17, 13, 420, 7), clr=tx_c)

            elif layout == 'stats':
                rect(sl, 0, 0, 13.333, 0.13, ac_c)
                txt(sl, 0.5, 0.28, 12.3, 0.85, title, sz=fit_sz(title, 28, 22, 90, 2), bold=True, clr=ti_c)
                stats = item.get('stats', [])[:3]
                n_st = len(stats)
                bw = (12.4 - (n_st - 1) * 0.4) / max(n_st, 1) if n_st else 12.4
                for i, stat in enumerate(stats):
                    sx = 0.45 + i * (bw + 0.4)
                    rect(sl, sx, 1.4, bw, 2.65, ac_c, a=12)
                    rect(sl, sx, 1.4, bw, 0.14, ac_c if i % 2 == 0 else ac2_c)
                    txt(sl, sx + 0.1, 1.65, bw - 0.2, 1.3,
                        stat.get('value', ''), sz=36, bold=True,
                        clr=ac_c if i % 2 == 0 else ac2_c, align='center')
                    txt(sl, sx + 0.1, 2.95, bw - 0.2, 0.7,
                        stat.get('label', ''), sz=15, clr=tx_c, align='center')
                if item.get('content'):
                    txt(sl, 0.5, 4.3, 12.3, 2.5, item['content'], sz=fit_sz(item['content'], 19, 14, 360, 6), clr=tx_c)

            elif layout == 'quote':
                txt(sl, 0.55, 0.4, 2.5, 2.5, '❝', sz=110, clr=ac_c)
                txt(sl, 1.5, 2.2, 10.6, 3.2, item.get('quote', ''),
                    sz=fit_sz(item.get('quote', ''), 26, 18, 340, 7), clr=ti_c, italic=True)
                rect(sl, 1.5, 5.6, 3.5, 0.07, ac_c)
                if title and title != f'Слайд {idx}':
                    txt(sl, 1.5, 5.85, 8.5, 0.6, f'— {title}', sz=15, clr=tx_c)

            elif layout == 'conclusion':
                rect(sl, 0, 6.7, 13.333, 0.8, ac2_c, a=85)
                rect(sl, 0, 6.7, 13.333, 0.8, ac_c,  a=25)
                rect(sl, 0, 0, 0.3, 6.7, ac2_c, a=90)
                rect(sl, 0.55, 1.2, 12.3, 4.5, ac2_c, a=6)
                txt(sl, 0.9, 1.45, 12.0, 1.4, title, sz=fit_sz(title, 40, 30, 90, 2), bold=True, clr=ti_c)
                rect(sl, 0.9, 2.9, 3.5, 0.07, ac2_c)
                content = item.get('content', item.get('subtitle', ''))
                if content:
                    txt(sl, 0.9, 3.15, 12.0, 3.0, content, sz=fit_sz(content, 22, 15, 420, 7), clr=tx_c)

            else:  # Обычный контентный слайд со списком тезисов.
                rect(sl, 0, 0, 0.23, 7.5, ac_c)
                txt(sl, 0.5, 0.28, 12.5, 1.0, title, sz=fit_sz(title, 32, 24, 90, 2), bold=True, clr=ti_c)
                rect(sl, 0.5, 1.3, 4.8, 0.06, ac_c, a=65)
                has_img = add_slide_image(sl, item, 8.3, 1.55, 4.65, 3.45, panel=True)
                image_placed = has_img
                bw = 7.5 if has_img else 12.6
                bullets = item.get('bullets') or [item.get('content', '')]
                bullet_text = '\n'.join(f'• {b}' for b in bullets if b)
                txt(sl, 0.5, 1.55, bw, 5.5,
                    bullet_text, sz=fit_sz(bullet_text, 20, 14, 520 if not has_img else 360, 8), clr=tx_c)

            if item.get('imageData') and not image_placed:
                # Запасное место для картинок в title/two_column/stats/quote/conclusion layout.
                image_placed = add_slide_image(sl, item, 9.0, 5.05, 3.7, 1.75, panel=True)

            slide_num(sl, idx)

        prs.save(output_path)
        return True
    except Exception as e:
        print(f'PPTX generation failed: {e}')
        traceback.print_exc()
        return False

def build_preview(slides: list[dict]) -> list[dict]:
    preview = []
    for i, s in enumerate(slides):
        preview.append({
            "index": i + 1,
            "layout": s.get("layout", "content"),
            "title": s.get("title", ""),
            "subtitle": s.get("subtitle", s.get("quote", "")),
            "content": s.get("content", ""),
            "quote": s.get("quote", ""),
            "bullets": s.get("bullets", []),
            "leftTitle": s.get("leftTitle", ""),
            "leftContent": s.get("leftContent", []),
            "rightTitle": s.get("rightTitle", ""),
            "rightContent": s.get("rightContent", []),
            "stats": s.get("stats", []),
            "image_prompt": s.get("image_prompt"),
            "imageData": _normalize_image_data_url(s.get("imageData")) if s.get("imageData") else "",
            "hasImage": bool(s.get("imageData")),
            "density": _slide_density(s),
            "plain_text": _slide_plain_text(s),
        })
    return preview

def extract_document_text(filename: str, content: bytes) -> str:
    fname = (filename or "").lower()
    if not content:
        return ""
    if fname.endswith(".pdf"):
        return extract_text_from_pdf(content)
    if fname.endswith(".docx"):
        return extract_text_from_docx(content)
    return ""

def create_presentation_from_data(
    prompt: str,
    slide_count: int,
    style: str,
    tone: str,
    generate_images: bool,
    rt_token: str,
    rt_service: str,
    document_filename: str = "",
    document_content: bytes = b"",
    progress_callback: Optional[ProgressCallback] = None,
) -> dict:
    """Полный pipeline генерации: документ -> структура -> QA -> картинки -> PPTX."""
    def progress(percent: int, title: str, detail: str = "") -> None:
        if progress_callback:
            progress_callback(percent, title, detail)

    style = normalize_style(style)
    session_id = str(uuid.uuid4())
    output_path = str(WORK_DIR / f"{session_id}.pptx")

    progress(3, "Чтение входных данных", "Проверяем параметры и загруженный документ.")
    doc_text = extract_document_text(document_filename, document_content)
    if doc_text:
        progress(12, "Документ прочитан", "Текст PDF/DOCX извлечён и добавлен в контекст.")
    else:
        progress(12, "Документ не загружен", "Будем строить презентацию по промпту пользователя.")

    slide_count = max(3, min(20, slide_count))
    progress(22, "Генерация структуры", "LLM формирует заголовки, layout-ы и тезисы слайдов.")
    structure = generate_slide_structure(prompt, doc_text, slide_count, style, tone, rt_token)
    progress(46, "Структура готова", "Запускаем QA-проверку и выравниваем плотность текста.")
    structure = review_and_refine_slide_structure(
        structure, prompt, doc_text, slide_count, style, tone, rt_token
    )
    progress(60, "QA-проверка завершена", "Слайды нормализованы, готовим визуальную часть.")
    structure["_presentai_context"] = {
        # Этот блок сохраняется в JSON, чтобы результат можно было редактировать позже.
        "prompt": prompt,
        "style": style,
        "tone": tone,
        "slide_count": slide_count,
        "document_filename": document_filename,
        "document_excerpt": doc_text[:7000],
        "rt_service": rt_service,
    }
    SESSION_SECRETS[session_id] = {
        # RT-токен нужен для точечного редактирования, но не должен попадать на диск.
        "rt_token": rt_token.strip(),
        "prompt": prompt,
        "style": style,
        "tone": tone,
        "slide_count": slide_count,
        "document_excerpt": doc_text[:7000],
        "rt_service": rt_service,
        "created_at": time.time(),
    }

    slides = structure.get("slides", [])
    image_warnings: list[str] = []
    if generate_images and rt_token.strip():
        progress(64, "Генерация изображений", "Запрашиваем картинки через выбранный RT image API.")
        image_started_at = time.time()
        image_attempts = 0
        image_successes = 0
        image_slots = max(1, min(RT_IMAGE_MAX_IMAGES, len(slides)))
        # Circuit breaker: счётчики последовательных ошибок по бэкенду.
        _cb_fails: dict[str, int] = {}
        _cb_disabled: set[str] = set()
        _active_service = rt_service

        def _cb_record_fail(svc: str) -> None:
            _cb_fails[svc] = _cb_fails.get(svc, 0) + 1
            if _cb_fails[svc] >= RT_IMAGE_CIRCUIT_BREAKER_THRESHOLD:
                _cb_disabled.add(svc)

        def _cb_record_ok(svc: str) -> None:
            _cb_fails[svc] = 0
            _cb_disabled.discard(svc)

        for idx, slide in enumerate(slides, start=1):
            elapsed = time.time() - image_started_at
            if elapsed >= RT_IMAGE_BATCH_TIMEOUT:
                image_warnings.append(
                    f"Генерация изображений остановлена после {int(elapsed)} секунд, чтобы не задерживать PPTX."
                )
                break
            if image_attempts >= RT_IMAGE_MAX_IMAGES:
                image_warnings.append(
                    f"Обработаны первые {RT_IMAGE_MAX_IMAGES} запросов изображений, успешно добавлено: {image_successes}. Остальные пропущены, чтобы презентация не зависла."
                )
                break

            # Если текущий сервис отключён circuit breaker-ом, попробуем альтернативный.
            if _active_service in _cb_disabled:
                alt = "sd" if _active_service == "yaArt" else "yaArt"
                if alt not in _cb_disabled:
                    _active_service = alt
                    image_warnings.append(f"Переключились на {alt}: {rt_service} временно недоступен.")
                else:
                    image_warnings.append("Оба image-бэкенда недоступны, пропускаем оставшиеся картинки.")
                    break

            # Первый и последний слайды — титульный/финальный, изображения не нужны.
            if idx == 1 or idx == len(slides):
                slide.pop("imageData", None)
                continue

            # Некоторые LLM-ответы не содержат image_prompt даже при включённых картинках.
            # Создаём запасной prompt, чтобы генерация действительно запускалась.
            _ensure_slide_image_prompt(slide, prompt)
            if slide.get("image_prompt"):
                image_attempts += 1
                image_progress = 64 + int((min(image_attempts - 1, image_slots) / image_slots) * 18)
                progress(
                    image_progress,
                    f"Генерация изображения {image_attempts}/{image_slots}",
                    f"Слайд {idx}: отправляем prompt в {_active_service}.",
                )
                ok, message = _generate_slide_image_if_needed(
                    slide,
                    rt_token.strip(),
                    _active_service,
                    topic=prompt,
                    instruction="",
                    force=False,
                )
                if ok and slide.get("imageData"):
                    image_successes += 1
                    _cb_record_ok(_active_service)
                else:
                    _cb_record_fail(_active_service)
                    if not slide.get("imageData"):
                        image_warnings.append(f"Слайд {idx}: {message}")
                progress(
                    64 + int((min(image_attempts, image_slots) / image_slots) * 18),
                    f"Изображение {image_attempts}/{image_slots} обработано",
                    message,
                )
                inter_delay = RT_IMAGE_RATELIMIT_PAUSE if (not ok and "429" in (message or "")) else 0.5
                time.sleep(inter_delay)
    elif generate_images:
        progress(72, "Изображения пропущены", "Генерация включена, но RT-токен не задан.")
        image_warnings.append("Генерация изображений включена, но RT-токен не задан.")
        for slide in slides:
            slide.pop("imageData", None)
    else:
        progress(72, "Изображения отключены", "Переходим к сборке PPTX без генерации картинок.")

    if image_warnings:
        structure["image_generation_warning"] = " ".join(image_warnings[:8])

    structure["slides"] = slides
    progress(88, "Сборка PPTX", "Создаём слайды, применяем тему и встраиваем изображения.")
    if not build_pptx(structure, output_path, style):
        raise HTTPException(500, "Ошибка генерации PPTX")

    progress(96, "Сохранение результата", "Записываем PPTX и JSON-сессию для редактора.")
    with open(WORK_DIR / f"{session_id}.json", "w", encoding="utf-8") as f:
        json.dump(structure, f, ensure_ascii=False, indent=2)

    return {
        "session_id": session_id,
        "title": structure.get("presentation_title", "Презентация"),
        "slide_count": len(slides),
        "preview": build_preview(slides),
        "warning": " ".join(x for x in [structure.get("generation_warning", ""), structure.get("review_warning", ""), structure.get("image_generation_warning", "")] if x).strip(),
        "source": structure.get("generation_source", ""),
        "review_source": structure.get("review_source", ""),
        "quality_review": structure.get("quality_review", {}),
        "style": style,
    }

async def create_presentation(
    prompt: str,
    slide_count: int,
    style: str,
    tone: str,
    generate_images: bool,
    rt_token: str,
    rt_service: str,
    document: Optional[UploadFile],
) -> dict:
    filename = ""
    content = b""
    if document and document.filename:
        filename = document.filename
        content = await document.read()
    return create_presentation_from_data(
        prompt, slide_count, style, tone, generate_images, rt_token, rt_service, filename, content
    )

def render_result_page(data: dict) -> str:
    safe_session = _xml_text(data.get("session_id", ""))
    safe_title = _xml_text(data.get("title", "Презентация готова"))
    slides_json = json.dumps(data.get("preview", []), ensure_ascii=False).replace("</", "<\\/")
    review_json = json.dumps(data.get("quality_review", {}), ensure_ascii=False).replace("</", "<\\/")
    warning = data.get("warning", "") or ""
    warning_json = json.dumps(warning, ensure_ascii=False).replace("</", "<\\/")
    source = data.get("source", "") or "неизвестен"
    review_source = data.get("review_source", "") or "QA-контроль"
    style = normalize_style(data.get("style", DEFAULT_STYLE))

    html = """<!doctype html>
<html lang="ru">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Презентация готова</title>
<style>
* { box-sizing: border-box; }
:root {
  --bg:#07040f; --panel:#120a26; --panel2:#1a1035; --text:#f8f5ff; --muted:#a99ac8;
  --accent:#8b5cf6; --accent2:#ec4e99; --ok:#22d3ee; --warn:#fed7aa; --border:rgba(139,92,246,.35);
}
body {
  margin:0; min-height:100vh; background:
  radial-gradient(circle at 10% 0%, rgba(139,92,246,.28), transparent 32%),
  radial-gradient(circle at 95% 20%, rgba(236,78,153,.18), transparent 35%),
  var(--bg);
  color:var(--text); font-family:Arial, sans-serif; padding:28px;
}
.wrap { max-width:1180px; margin:0 auto; }
.top { display:flex; gap:16px; align-items:flex-start; justify-content:space-between; margin-bottom:18px; }
h1 { margin:0 0 8px; font-size:clamp(28px,4vw,44px); line-height:1.06; }
.meta { color:var(--muted); line-height:1.45; font-size:14px; }
.actions { display:flex; gap:10px; flex-wrap:wrap; justify-content:flex-end; }
.btn, button {
  display:inline-flex; align-items:center; justify-content:center; gap:8px;
  border:1px solid var(--border); background:linear-gradient(135deg,var(--accent),var(--accent2));
  color:white; padding:12px 16px; border-radius:12px; text-decoration:none; font-weight:800;
  cursor:pointer; transition:.2s; font-size:14px;
}
.btn:hover, button:hover { transform:translateY(-1px); filter:brightness(1.06); }
button.secondary, .btn.secondary { background:rgba(255,255,255,.06); }
button:disabled { opacity:.55; cursor:wait; transform:none; }
.notice {
  margin:12px 0; padding:12px 14px; border-radius:14px; border:1px solid rgba(249,115,22,.35);
  background:rgba(249,115,22,.12); color:var(--warn); line-height:1.45; display:none;
}
.qa {
  display:grid; grid-template-columns: minmax(0,1fr) auto; gap:12px; align-items:center;
  margin:14px 0 18px; padding:14px; border:1px solid var(--border); border-radius:16px;
  background:rgba(18,10,38,.78);
}
.qa strong { color:var(--ok); }
.score { min-width:88px; height:88px; border-radius:20px; display:grid; place-items:center;
  background:rgba(139,92,246,.18); border:1px solid var(--border); font-size:28px; font-weight:900; }
.layout {
  display:grid; grid-template-columns: 1.35fr .9fr; gap:18px; align-items:start;
}
.card {
  background:rgba(18,10,38,.86); border:1px solid var(--border); border-radius:18px; padding:16px;
  box-shadow:0 24px 80px rgba(0,0,0,.22);
}
.preview-toolbar { display:flex; justify-content:space-between; align-items:center; gap:10px; color:var(--muted); margin-bottom:12px; font-size:13px; }
.preview-slide {
  aspect-ratio:16/9; width:100%; border-radius:16px; overflow:hidden; padding:5.2%;
  background:linear-gradient(150deg,#0a0f1e,#1a0a3e); border:1px solid rgba(255,255,255,.08);
  box-shadow:inset 0 0 120px rgba(119,0,255,.2); display:flex; flex-direction:column; justify-content:center;
}
.preview-slide h2 { margin:0 0 18px; font-size:clamp(26px,4vw,48px); line-height:1.08; }
.preview-slide p, .preview-slide li { font-size:clamp(14px,1.7vw,22px); color:rgba(255,255,255,.84); line-height:1.42; }
.preview-slide ul { margin:0; padding-left:24px; display:grid; gap:8px; }
.preview-with-image { display:grid; grid-template-columns:minmax(0,1fr) 42%; gap:18px; align-items:center; width:100%; }
.preview-image { width:100%; max-height:100%; aspect-ratio:4/3; object-fit:contain; border-radius:14px; background:rgba(255,255,255,.06); border:1px solid rgba(255,255,255,.12); box-shadow:0 18px 42px rgba(0,0,0,.22); }
.preview-image-only { display:grid; gap:14px; align-items:center; justify-items:center; }
.preview-columns { display:grid; grid-template-columns:1fr 1fr; gap:16px; }
.preview-column { padding:14px; border-radius:12px; background:rgba(255,255,255,.07); }
.preview-column h3 { margin:0 0 10px; color:white; }
.preview-stats { display:grid; grid-template-columns:repeat(3,1fr); gap:12px; margin-bottom:16px; }
.preview-stat { background:rgba(255,255,255,.08); border-radius:14px; padding:14px; text-align:center; }
.preview-stat strong { display:block; font-size:28px; color:#c4b5fd; margin-bottom:6px; }
.editor { display:grid; gap:10px; }
label { font-size:12px; color:var(--muted); font-weight:800; text-transform:uppercase; letter-spacing:.04em; }
input, textarea {
  width:100%; border:1px solid var(--border); background:#0e0820; color:var(--text);
  border-radius:12px; padding:11px 12px; font:14px Arial, sans-serif; outline:none;
}
textarea { min-height:130px; resize:vertical; line-height:1.45; }
#edit-instruction { min-height:92px; }
.help { color:var(--muted); font-size:12px; line-height:1.45; }
.status { min-height:22px; color:var(--ok); font-size:13px; line-height:1.45; }
.slides-grid { display:grid; grid-template-columns:repeat(auto-fill,minmax(190px,1fr)); gap:12px; margin-top:18px; }
.slide-card { border:1px solid rgba(139,92,246,.22); background:#100821; border-radius:14px; overflow:hidden; cursor:pointer; transition:.2s; }
.slide-card:hover { transform:translateY(-2px); border-color:var(--border); }
.slide-card.active { border-color:var(--accent); box-shadow:0 0 0 2px rgba(139,92,246,.28); }
.thumb { min-height:126px; padding:14px; background:linear-gradient(150deg,#0d1117,#1a1035); position:relative; }
.num { position:absolute; top:8px; right:10px; font-size:11px; color:rgba(255,255,255,.38); font-weight:900; }
.badge { font-size:10px; color:#c4b5fd; font-weight:900; text-transform:uppercase; margin-bottom:8px; letter-spacing:.05em; }
.thumb-title { font-weight:900; font-size:13px; line-height:1.25; display:-webkit-box; -webkit-line-clamp:2; -webkit-box-orient:vertical; overflow:hidden; }
.thumb-text { margin-top:8px; color:rgba(255,255,255,.54); font-size:10px; line-height:1.35; display:-webkit-box; -webkit-line-clamp:4; -webkit-box-orient:vertical; overflow:hidden; }
.slide-foot { padding:9px 12px; display:flex; justify-content:space-between; gap:8px; color:var(--muted); font-size:11px; border-top:1px solid rgba(139,92,246,.16); }
.density.ok { color:#86efac; } .density.mid { color:#fed7aa; } .density.bad { color:#fca5a5; }
.issue-list { margin:10px 0 0; padding-left:18px; color:var(--muted); font-size:13px; line-height:1.45; }
@media(max-width:900px){ body{padding:18px}.layout{grid-template-columns:1fr}.top{display:block}.actions{justify-content:flex-start;margin-top:14px}.qa{grid-template-columns:1fr}.score{height:64px;min-width:64px;font-size:22px}.preview-columns,.preview-stats{grid-template-columns:1fr} }
</style>
</head>
<body>
<div class="wrap">
  <div class="top">
    <div>
      <h1>__TITLE__</h1>
      <div class="meta">
        Источник генерации: <strong>__SOURCE__</strong><br>
        Проверка второй моделью: <strong>__REVIEW_SOURCE__</strong>
      </div>
    </div>
    <div class="actions">
      <a class="btn" id="download-btn" href="/api/download/__SESSION_ID__">⬇ Скачать текущий PPTX</a>
      <a class="btn secondary" href="/">Создать ещё</a>
    </div>
  </div>

  <div class="notice" id="warning">__WARNING__</div>

  <div class="qa" id="qa-panel">
    <div>
      <strong>QA-отчёт</strong>
      <div class="meta" id="qa-summary">Проверка структуры, достоверности, плотности текста и визуального баланса.</div>
      <ul class="issue-list" id="qa-issues"></ul>
    </div>
    <div class="score" id="qa-score">—</div>
  </div>

  <div class="layout">
    <div class="card">
      <div class="preview-toolbar">
        <span id="preview-counter">Слайд 1</span>
        <span>Предпросмотр выбранного слайда</span>
      </div>
      <div class="preview-slide" id="preview-slide"></div>
    </div>

    <div class="card editor">
      <label>Заголовок выбранного слайда</label>
      <input id="edit-title" placeholder="Заголовок">
      <label>Текст / пункты / цитата</label>
      <textarea id="edit-body" placeholder="Каждый пункт — с новой строки"></textarea>
      <label>Инструкция второй LLM для точечной правки</label>
      <textarea id="edit-instruction" placeholder="Например: сократи текст, сделай слайд красивее, проверь цифры по документу, поправь структуру"></textarea>
      <label>RT-токен для правки, если сервер потерял исходный токен</label>
      <input id="edit-token" placeholder="Оставьте пустым, если токен вводили при генерации">
      <div class="help">
        «Сохранить вручную» пересобирает PPTX только с вашими изменениями выбранного слайда.
        «Проверить и исправить LLM» отправляет во вторую модель только выбранный слайд, без регенерации всей презентации.
      </div>
      <button class="secondary" id="save-manual" onclick="saveSelectedSlide(false)">💾 Сохранить вручную</button>
      <button id="save-llm" onclick="saveSelectedSlide(true)">🧠 Проверить и исправить LLM</button>
      <div class="status" id="edit-status"></div>
    </div>
  </div>

  <div class="slides-grid" id="slides-grid"></div>
</div>

<script>
const sessionId = "__SESSION_ID__";
let currentStyle = "__STYLE__";
let currentSlides = __SLIDES_JSON__;
let qualityReview = __REVIEW_JSON__;
let activeIndex = 0;

const labels = {
  title:"Титульный", content:"Контент", two_column:"Два столбца",
  stats:"Статистика", quote:"Цитата", section_break:"Раздел",
  conclusion:"Заключение", image:"Изображение"
};

function escapeHtml(value) {
  return String(value ?? '').replace(/[&<>"']/g, ch => ({
    '&':'&amp;', '<':'&lt;', '>':'&gt;', '"':'&quot;', "'":'&#39;'
  }[ch]));
}

function bodyFromSlide(slide) {
  if (!slide) return '';
  if (slide.layout === 'title') return slide.subtitle || '';
  if (slide.layout === 'quote') return slide.quote || slide.content || '';
  if (slide.layout === 'conclusion' || slide.layout === 'section_break') return slide.content || slide.subtitle || '';
  if (slide.layout === 'stats') return slide.content || '';
  if (slide.layout === 'two_column') return [...(slide.leftContent || []), ...(slide.rightContent || [])].join('\\n');
  if (slide.bullets && slide.bullets.length) return slide.bullets.join('\\n');
  return slide.content || slide.subtitle || slide.quote || '';
}

function slideContentHtml(s) {
  if (!s) return '';
  if (s.stats && s.stats.length) {
    return `<div class="preview-stats">${s.stats.slice(0,3).map(st =>
      `<div class="preview-stat"><strong>${escapeHtml(st.value)}</strong><span>${escapeHtml(st.label)}</span></div>`
    ).join('')}</div><p>${escapeHtml(s.content || '')}</p>`;
  }
  if (s.layout === 'two_column') {
    return `<div class="preview-columns">
      <div class="preview-column"><h3>${escapeHtml(s.leftTitle || 'Блок 1')}</h3><ul>${(s.leftContent || []).map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul></div>
      <div class="preview-column"><h3>${escapeHtml(s.rightTitle || 'Блок 2')}</h3><ul>${(s.rightContent || []).map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul></div>
    </div>`;
  }
  if (s.bullets && s.bullets.length) {
    return `<ul>${s.bullets.map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul>`;
  }
  return `<p>${escapeHtml(s.quote || s.content || s.subtitle || '')}</p>`;
}

function slideImageHtml(s) {
  if (!s || !s.imageData) return '';
  return `<img class="preview-image" src="${escapeHtml(s.imageData)}" alt="Изображение слайда">`;
}

function slideBodyHtml(s) {
  const body = slideContentHtml(s);
  const image = slideImageHtml(s);
  if (!image) return body;
  if (!body || body === '<p></p>') return `<div class="preview-image-only">${image}</div>`;
  return `<div class="preview-with-image"><div>${body}</div>${image}</div>`;
}
function densityClass(slide) {
  const ratio = Number(slide?.density?.ratio || 0);
  if (ratio > 1.12) return 'bad';
  if (ratio > 0.86) return 'mid';
  return 'ok';
}

function renderQa() {
  const score = qualityReview?.score ?? '—';
  document.getElementById('qa-score').textContent = score;
  document.getElementById('qa-summary').textContent = qualityReview?.summary || 'Проверка выполнена.';
  const issues = Array.isArray(qualityReview?.issues) ? qualityReview.issues.slice(0, 8) : [];
  document.getElementById('qa-issues').innerHTML = issues.length
    ? issues.map(x => `<li>Слайд ${escapeHtml(x.slide || '?')}: ${escapeHtml(x.issue || x)}</li>`).join('')
    : '<li>Критичных проблем не найдено.</li>';
}

function renderCards() {
  const grid = document.getElementById('slides-grid');
  grid.innerHTML = '';
  currentSlides.forEach((slide, idx) => {
    const card = document.createElement('div');
    card.className = 'slide-card' + (idx === activeIndex ? ' active' : '');
    card.onclick = () => selectSlide(idx);
    const body = slide.plain_text || bodyFromSlide(slide);
    const ratio = slide?.density?.ratio ?? 0;
    card.innerHTML = `
      <div class="thumb">
        <div class="num">${idx + 1}</div>
        <div class="badge">${labels[slide.layout] || slide.layout || 'Слайд'}</div>
        <div class="thumb-title">${escapeHtml(slide.title || ('Слайд ' + (idx + 1)))}</div>
        <div class="thumb-text">${escapeHtml(body)}</div>
      </div>
      <div class="slide-foot">
        <span>${slide.hasImage ? '🖼 image' : 'text'}</span>
        <span class="density ${densityClass(slide)}">density ${escapeHtml(ratio)}</span>
      </div>`;
    grid.appendChild(card);
  });
}

function selectSlide(index) {
  activeIndex = Math.max(0, Math.min(currentSlides.length - 1, Number(index) || 0));
  const slide = currentSlides[activeIndex];
  document.getElementById('preview-counter').textContent = `Слайд ${activeIndex + 1} из ${currentSlides.length}`;
  document.getElementById('preview-slide').innerHTML = `<h2>${escapeHtml(slide.title || ('Слайд ' + (activeIndex + 1)))}</h2>${slideBodyHtml(slide)}`;
  document.getElementById('edit-title').value = slide.title || '';
  document.getElementById('edit-body').value = bodyFromSlide(slide);
  document.getElementById('edit-instruction').value = '';
  renderCards();
}

async function saveSelectedSlide(useLLM) {
  const status = document.getElementById('edit-status');
  const btnManual = document.getElementById('save-manual');
  const btnLlm = document.getElementById('save-llm');
  const selected = currentSlides[activeIndex];
  if (!selected) return;
  btnManual.disabled = true;
  btnLlm.disabled = true;
  status.textContent = useLLM ? 'Отправляем выбранный слайд во вторую LLM...' : 'Сохраняем выбранный слайд и пересобираем PPTX...';

  try {
    const resp = await fetch(`/api/session/${encodeURIComponent(sessionId)}/slide/${activeIndex + 1}/edit`, {
      method: 'POST',
      headers: {'Content-Type':'application/json'},
      body: JSON.stringify({
        title: document.getElementById('edit-title').value,
        body: document.getElementById('edit-body').value,
        instruction: useLLM ? document.getElementById('edit-instruction').value : '',
        use_llm: !!useLLM,
        rt_token: document.getElementById('edit-token').value
      })
    });
    const data = await resp.json().catch(() => ({}));
    if (!resp.ok) throw new Error(data.detail || 'Не удалось обновить слайд');
    currentSlides = data.preview || currentSlides;
    qualityReview = data.quality_review || qualityReview;
    document.getElementById('download-btn').href = `/api/download/${sessionId}?t=${Date.now()}`;
    renderQa();
    selectSlide(Math.min(activeIndex, currentSlides.length - 1));
    const report = data.slide_review?.source ? ` Источник: ${data.slide_review.source}.` : '';
    status.textContent = 'Готово: обновлён только выбранный слайд, PPTX пересобран.' + report;
  } catch (err) {
    status.textContent = 'Ошибка: ' + (err.message || String(err));
  } finally {
    btnManual.disabled = false;
    btnLlm.disabled = false;
  }
}

(function init() {
  const warning = __WARNING_JSON__;
  if (warning.trim()) {
    const node = document.getElementById('warning');
    node.style.display = 'block';
    node.textContent = warning;
  }
  renderQa();
  renderCards();
  selectSlide(0);
})();
</script>
</body>
</html>"""
    return (
        html
        .replace("__SESSION_ID__", safe_session)
        .replace("__TITLE__", safe_title)
        .replace("__SOURCE__", _xml_text(source))
        .replace("__REVIEW_SOURCE__", _xml_text(review_source))
        .replace("__WARNING__", _xml_text(warning))
        .replace("__WARNING_JSON__", warning_json)
        .replace("__STYLE__", _xml_text(style))
        .replace("__SLIDES_JSON__", slides_json)
        .replace("__REVIEW_JSON__", review_json)
    )

def render_loading_page(job_id: str) -> str:
    safe_job_id = _xml_text(job_id)
    return f"""<!doctype html>
<html lang="ru"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Генерация презентации</title>
<style>
*{{box-sizing:border-box}}body{{margin:0;min-height:100vh;background:#07040f;color:#f5f1ff;font-family:Arial,sans-serif;display:grid;place-items:center;padding:28px}}
.card{{width:min(760px,100%);background:#100821;border:1px solid rgba(139,92,246,.35);border-radius:22px;padding:38px;text-align:center;box-shadow:0 24px 90px rgba(139,92,246,.22)}}
.logo{{width:58px;height:58px;border-radius:16px;margin:0 auto 18px;display:grid;place-items:center;background:linear-gradient(135deg,#8b5cf6,#ec4e99);font-size:28px}}
h1{{margin:0 0 8px;font-size:clamp(28px,5vw,44px)}}p{{margin:0 0 24px;color:#a99ac8;line-height:1.5}}
.percent{{font-size:52px;font-weight:800;color:#a78bfa;margin:4px 0 16px}}
.bar{{height:14px;background:#21143e;border-radius:999px;overflow:hidden;border:1px solid rgba(255,255,255,.08)}}
.fill{{height:100%;width:0;background:linear-gradient(90deg,#8b5cf6,#ec4e99);border-radius:999px;transition:width .45s ease}}
.steps{{display:grid;gap:10px;margin-top:24px;text-align:left}}.step{{padding:12px 14px;border-radius:12px;background:rgba(255,255,255,.04);color:#b9aad6}}.step.done{{color:#fff;background:rgba(34,211,238,.12)}}
.meta{{margin-top:18px;color:#7f6aa8;font-size:13px}}
.error{{display:none;margin-top:20px;padding:14px;border-radius:12px;background:rgba(239,68,68,.14);color:#fecaca;text-align:left;white-space:pre-wrap}}
a{{color:#c4b5fd}}
</style></head><body>
<main class="card">
  <div class="logo">✦</div>
  <h1 id="title">Генерируем презентацию</h1>
  <p id="sub">Не закрывайте страницу. Сервис анализирует запрос и собирает PPTX.</p>
  <div class="percent" id="percent">0%</div>
  <div class="bar"><div class="fill" id="fill"></div></div>
  <div class="steps">
    <div class="step" id="s1">Чтение документа и параметров</div>
    <div class="step" id="s2">Генерация структуры и QA-проверка второй LLM</div>
    <div class="step" id="s3">Правка структуры, текста и визуального баланса</div>
    <div class="step" id="s4">Сборка PPTX и страницы результата</div>
  </div>
  <div class="meta" id="meta">Ожидаем запуск задачи...</div>
  <div class="error" id="error"></div>
</main>
<script>
const jobId = "{safe_job_id}";
const percent = document.getElementById('percent');
const fill = document.getElementById('fill');
const title = document.getElementById('title');
const sub = document.getElementById('sub');
const meta = document.getElementById('meta');
const error = document.getElementById('error');
let value = 0;
let finished = false;
function setProgress(next, nextTitle, nextDetail) {{
  value = Math.max(value, Math.min(100, Math.round(next)));
  percent.textContent = value + '%';
  fill.style.width = value + '%';
  if (nextTitle) title.textContent = nextTitle;
  if (nextDetail) sub.textContent = nextDetail;
  if (value >= 18) document.getElementById('s1').classList.add('done');
  if (value >= 42) document.getElementById('s2').classList.add('done');
  if (value >= 70) document.getElementById('s3').classList.add('done');
  if (value >= 92) document.getElementById('s4').classList.add('done');
}}
function showError(message) {{
  finished = true;
  error.style.display = 'block';
  error.innerHTML = '<strong>Ошибка:</strong> ' + String(message || 'Неизвестная ошибка') + '<br><br><a href="/">Вернуться к форме</a>';
}}
async function pollStatus() {{
  if (finished) return;
  try {{
    const resp = await fetch('/api/jobs/' + encodeURIComponent(jobId) + '/status', {{cache:'no-store'}});
    const data = await resp.json().catch(() => ({{}}));
    if (!resp.ok) throw new Error(data.detail || 'Не удалось получить статус');
    setProgress(data.progress || 0, data.title, data.detail);
    meta.textContent = data.status === 'running'
      ? 'Задача выполняется на сервере. Проценты обновляются по реальным этапам.'
      : data.status === 'queued'
        ? 'Задача в очереди запуска.'
        : '';
    if (data.status === 'done' && data.session_id) {{
      finished = true;
      setProgress(100, 'Презентация готова', 'Открываем страницу скачивания и предпросмотра.');
      setTimeout(() => {{ window.location.href = '/result/' + encodeURIComponent(data.session_id); }}, 700);
    }}
    if (data.status === 'error') {{
      showError(data.error || 'Ошибка генерации');
    }}
  }} catch (err) {{
    meta.textContent = 'Ждём ответ сервера статуса...';
  }}
}}
async function run() {{
  const pollTimer = setInterval(pollStatus, 700);
  try {{
    setProgress(1, 'Запускаем задачу', 'Передаём параметры генерации на сервер.');
    await pollStatus();
    const controller = new AbortController();
    const timeoutId = setTimeout(() => controller.abort(), 360000);
    const resp = await fetch('/api/jobs/' + encodeURIComponent(jobId) + '/run', {{
      method:'POST',
      signal: controller.signal
    }});
    clearTimeout(timeoutId);
    const data = await resp.json().catch(() => ({{}}));
    if (!resp.ok) throw new Error(data.detail || 'Ошибка генерации');
    if (data.session_id) {{
      finished = true;
      clearInterval(pollTimer);
      setProgress(100, 'Презентация готова', 'Открываем страницу скачивания и предпросмотра.');
      setTimeout(() => {{ window.location.href = '/result/' + encodeURIComponent(data.session_id); }}, 700);
    }}
  }} catch (err) {{
    clearInterval(pollTimer);
    const message = err.name === 'AbortError'
      ? 'Сервер не ответил за 6 минут. Попробуйте выключить генерацию изображений или уменьшить количество слайдов.'
      : String(err.message || err);
    showError(message);
  }}
}}
run();
</script>
</body></html>"""

# ── API-маршруты ──────────────────────────────────────────────────────────────

@app.post("/api/generate")
async def generate_presentation(
    prompt: str = Form(...),
    slide_count: int = Form(8),
    style: str = Form(DEFAULT_STYLE),
    tone: str = Form("professional"),
    generate_images: bool = Form(False),
    rt_token: str = Form(""),
    rt_service: str = Form("sd"),
    document: Optional[UploadFile] = File(None),
):
    try:
        return JSONResponse(await create_presentation(
            prompt, slide_count, style, tone, generate_images, rt_token, rt_service, document
        ))

    except json.JSONDecodeError as e:
        raise HTTPException(500, f"Ошибка парсинга ответа LLM: {e}")
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))


@app.post("/loading", response_class=HTMLResponse)
async def create_loading_job(
    prompt: str = Form(...),
    slide_count: int = Form(8),
    style: str = Form(DEFAULT_STYLE),
    tone: str = Form("professional"),
    generate_images: bool = Form(False),
    rt_token: str = Form(""),
    rt_service: str = Form("sd"),
    document: Optional[UploadFile] = File(None),
):
    filename = ""
    content = b""
    if document and document.filename:
        filename = document.filename
        content = await document.read()

    job_id = str(uuid.uuid4())
    JOBS[job_id] = {
        "prompt": prompt,
        "slide_count": slide_count,
        "style": style,
        "tone": tone,
        "generate_images": generate_images,
        "rt_token": rt_token,
        "rt_service": rt_service,
        "document_filename": filename,
        "document_content": content,
        "status": "queued",
        "progress": 0,
        "title": "Ожидаем запуска",
        "detail": "Задача создана, страница загрузки запускает генерацию.",
        "result": None,
        "error": "",
        "created_at": time.time(),
        "updated_at": time.time(),
    }
    return render_loading_page(job_id)


@app.post("/api/jobs/{job_id}/run")
async def run_loading_job(job_id: str):
    job = JOBS.get(job_id)
    if not job:
        raise HTTPException(404, "Задача генерации не найдена. Вернитесь к форме и запустите генерацию заново.")
    if job.get("status") == "done" and isinstance(job.get("result"), dict):
        return JSONResponse(job["result"])
    if job.get("status") == "running":
        return JSONResponse(_job_public_status(job_id, job))
    try:
        _job_update(job_id, 1, "Запуск генерации", "Сервер начал обработку задачи.", status="running")

        def progress(percent: int, title: str, detail: str = "") -> None:
            _job_update(job_id, percent, title, detail, status="running")

        data = await asyncio.to_thread(
            create_presentation_from_data,
            job["prompt"],
            job["slide_count"],
            job["style"],
            job["tone"],
            job["generate_images"],
            job["rt_token"],
            job["rt_service"],
            job["document_filename"],
            job["document_content"],
            progress,
        )
        job.pop("document_content", None)
        job.pop("rt_token", None)
        _job_update(job_id, 100, "Презентация готова", "Файл собран, открываем результат.", status="done", result=data)
        return JSONResponse(data)
    except Exception as e:
        traceback.print_exc()
        job.pop("document_content", None)
        job.pop("rt_token", None)
        _job_update(job_id, None, "Ошибка генерации", str(e), status="error", error=str(e))
        raise HTTPException(500, str(e))


@app.get("/api/jobs/{job_id}/status")
async def loading_job_status(job_id: str):
    job = JOBS.get(job_id)
    if not job:
        raise HTTPException(404, "Задача генерации не найдена")
    return JSONResponse(_job_public_status(job_id, job))


@app.post("/generate", response_class=HTMLResponse)
async def generate_presentation_page(
    prompt: str = Form(...),
    slide_count: int = Form(8),
    style: str = Form(DEFAULT_STYLE),
    tone: str = Form("professional"),
    generate_images: bool = Form(False),
    rt_token: str = Form(""),
    rt_service: str = Form("sd"),
    document: Optional[UploadFile] = File(None),
):
    try:
        data = await create_presentation(
            prompt, slide_count, style, tone, generate_images, rt_token, rt_service, document
        )
        return render_result_page(data)
    except Exception as e:
        traceback.print_exc()
        return HTMLResponse(f"<h1>Ошибка генерации</h1><pre>{_xml_text(e)}</pre><p><a href='/'>Назад</a></p>", status_code=500)


@app.post("/api/rebuild/{session_id}")
async def rebuild_pptx(session_id: str, request: Request):
    """Пересобирает PPTX из отредактированных в браузере слайдов."""
    try:
        payload = await request.json()
        slides = payload.get("slides")
        existing = {}
        try:
            existing = load_session_structure(session_id)
        except HTTPException:
            existing = {}

        context = _session_context(existing, session_id) if existing else {}
        style = normalize_style(payload.get("style") or context.get("style") or DEFAULT_STYLE)
        title = payload.get("title") or existing.get("presentation_title") or "Презентация"
        if not isinstance(slides, list) or not slides:
            raise HTTPException(400, "Нет данных слайдов для пересборки")

        old_slides = existing.get("slides", []) if isinstance(existing.get("slides"), list) else []
        merged_slides = []
        for idx, slide in enumerate(slides):
            merged = copy.deepcopy(slide)
            if idx < len(old_slides):
                # Ручное редактирование текста не должно стирать уже вложенные картинки.
                old = old_slides[idx]
                if old.get("imageData") and not merged.get("imageData"):
                    merged["imageData"] = old.get("imageData")
                if old.get("image_prompt") and not merged.get("image_prompt"):
                    merged["image_prompt"] = old.get("image_prompt")
            cleaned, _, _ = _sanitize_slide_for_quality(
                merged, idx, len(slides), _doc_keyword_set(context.get("document_excerpt", "")), context.get("document_excerpt", "")
            )
            merged_slides.append(cleaned)

        structure = copy.deepcopy(existing) if existing else {}
        structure.update({
            "presentation_title": title,
            "metadata": {"title": title, "author": "PresentAI", "contact": ""},
            "slides": merged_slides,
        })
        structure.setdefault("_presentai_context", {})
        structure["_presentai_context"].update({
            "style": style,
            "slide_count": len(merged_slides),
            "document_excerpt": context.get("document_excerpt", ""),
            "prompt": context.get("prompt", ""),
            "tone": context.get("tone", "professional"),
        })

        save_session_structure(session_id, structure, style)

        return JSONResponse({
            "session_id": session_id,
            "title": title,
            "slide_count": len(merged_slides),
            "preview": build_preview(merged_slides),
            "quality_review": structure.get("quality_review", {}),
            "style": style,
        })
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))


@app.get("/api/session/{session_id}")
async def get_session(session_id: str):
    structure = load_session_structure(session_id)
    slides = structure.get("slides", []) or []
    context = _session_context(structure, session_id)
    return JSONResponse({
        "session_id": session_id,
        "title": structure.get("presentation_title", "Презентация"),
        "slide_count": len(slides),
        "preview": build_preview(slides),
        "source": structure.get("generation_source", ""),
        "warning": " ".join(x for x in [structure.get("generation_warning", ""), structure.get("review_warning", ""), structure.get("image_generation_warning", "")] if x).strip(),
        "review_source": structure.get("review_source", ""),
        "quality_review": structure.get("quality_review", {}),
        "style": normalize_style(context.get("style", DEFAULT_STYLE)),
    })


@app.post("/api/session/{session_id}/slide/{slide_index}/edit")
async def edit_session_slide(session_id: str, slide_index: int, request: Request):
    """Редактирует один выбранный слайд и не регенерирует всю презентацию."""
    try:
        payload = await request.json()
        structure = load_session_structure(session_id)
        slides = structure.get("slides", [])
        if not isinstance(slides, list) or not slides:
            raise HTTPException(400, "В сессии нет слайдов")
        idx = slide_index - 1
        if idx < 0 or idx >= len(slides):
            raise HTTPException(404, "Слайд не найден")

        context = _session_context(structure, session_id)
        style = normalize_style(payload.get("style") or context.get("style") or DEFAULT_STYLE)
        doc_text = context.get("document_excerpt", "")
        topic = context.get("prompt", "") or structure.get("presentation_title", "")
        rt_service = payload.get("rt_service") or context.get("rt_service") or "sd"
        runtime_token = SESSION_SECRETS.get(session_id, {}).get("rt_token", "")
        rt_token = (payload.get("rt_token") or "").strip() or runtime_token

        before_image_prompt = _clean_text(slides[idx].get("image_prompt"))
        before_had_image = bool(slides[idx].get("imageData"))

        # Сначала применяем ручные поля из редактора только к выбранному слайду.
        slides[idx] = _slide_from_editor_payload(slides[idx], payload, idx, len(slides))
        structure["slides"] = slides

        instruction = _clean_text(payload.get("instruction") or "")
        use_llm = bool(payload.get("use_llm"))
        if use_llm:
            structure, slide_report = edit_single_slide_with_reviewer(
                structure, idx, instruction, rt_token, doc_text
            )
        else:
            structure, slide_report = edit_single_slide_with_reviewer(
                structure, idx, "", "", doc_text
            )

        # Если точечная правка просит картинку, LLM обновляет только image_prompt.
        # Сам файл изображения генерируется отдельно и только для выбранного слайда.
        edited_slides = structure.get("slides", []) or []
        if idx < len(edited_slides):
            selected_slide = edited_slides[idx]
            after_image_prompt = _clean_text(selected_slide.get("image_prompt"))
            prompt_changed = bool(after_image_prompt and after_image_prompt != before_image_prompt)
            image_requested = _looks_like_image_request(instruction) or prompt_changed
            should_attach_image = bool(
                use_llm and (image_requested or (after_image_prompt and not selected_slide.get("imageData")))
            )
            if should_attach_image:
                if prompt_changed and before_had_image:
                    selected_slide.pop("imageData", None)
                if image_requested and not after_image_prompt:
                    _ensure_slide_image_prompt(selected_slide, topic, instruction, force=True)
                force_regenerate = prompt_changed or image_requested or not selected_slide.get("imageData")
                ok, image_message = _generate_slide_image_if_needed(
                    selected_slide,
                    rt_token,
                    rt_service,
                    topic=topic,
                    instruction=instruction,
                    force=force_regenerate,
                )
                if ok:
                    slide_report.setdefault("fixes", []).append(image_message)
                else:
                    slide_report.setdefault("issues", []).append(image_message)
                edited_slides[idx] = selected_slide
                structure["slides"] = edited_slides

        # Обновляем QA-метаданные, не трогая замечания по остальным слайдам.
        old_review = structure.get("quality_review") if isinstance(structure.get("quality_review"), dict) else {}
        old_issues = [x for x in (old_review.get("issues") or []) if not (isinstance(x, dict) and x.get("slide") == slide_index)]
        slide_issues = []
        for issue in slide_report.get("issues", []) or []:
            if isinstance(issue, dict):
                slide_issues.append(issue)
            else:
                slide_issues.append({"slide": slide_index, "severity": "low", "issue": str(issue), "fix": ""})
        slide_fixes = [str(x) for x in (slide_report.get("fixes") or [])]
        structure["quality_review"] = {
            "source": slide_report.get("source") or old_review.get("source") or "Редактор слайда",
            "score": old_review.get("score", 90),
            "summary": f"Слайд {slide_index} обновлён отдельно; остальная презентация не регенерировалась.",
            "issues": (old_issues + slide_issues)[:50],
            "corrections": (old_review.get("corrections") or [])[:35] + slide_fixes[:15],
        }
        structure.setdefault("_presentai_context", {})
        structure["_presentai_context"].update({
            "style": style,
            "document_excerpt": doc_text,
            "prompt": context.get("prompt", ""),
            "tone": context.get("tone", "professional"),
            "rt_service": rt_service,
            "slide_count": len(structure.get("slides", [])),
        })

        save_session_structure(session_id, structure, style)

        return JSONResponse({
            "session_id": session_id,
            "title": structure.get("presentation_title", "Презентация"),
            "slide_count": len(structure.get("slides", [])),
            "preview": build_preview(structure.get("slides", [])),
            "quality_review": structure.get("quality_review", {}),
            "slide_review": slide_report,
            "style": style,
        })
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(500, str(e))


@app.get("/api/demo")
async def demo_pptx():
    """Генерирует и скачивает демо-PPTX без JavaScript и API-ключей."""
    session_id = str(uuid.uuid4())
    output_path = str(WORK_DIR / f"{session_id}.pptx")
    structure = generate_slide_structure(
        "Демо-презентация: AI-генератор презентаций",
        "Проверка локальной генерации без API-ключей. Сервис принимает промпт, документ PDF или DOCX, настройки стиля и тона, затем собирает PPTX.",
        6,
        DEFAULT_STYLE,
        "professional",
        "",
    )
    if not build_pptx(structure, output_path, DEFAULT_STYLE):
        raise HTTPException(500, "Ошибка генерации demo PPTX")
    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentai-demo.pptx",
    )


@app.get("/api/download/{session_id}")
async def download_pptx(session_id: str):
    path = WORK_DIR / f"{session_id}.pptx"
    if not path.exists():
        raise HTTPException(404, "Файл не найден")
    return FileResponse(
        str(path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentation.pptx",
    )


@app.get("/result/{session_id}", response_class=HTMLResponse)
async def result_page(session_id: str):
    json_path = WORK_DIR / f"{session_id}.json"
    pptx_path = WORK_DIR / f"{session_id}.pptx"
    if not json_path.exists() or not pptx_path.exists():
        return HTMLResponse(
            "<h1>Презентация не найдена</h1><p><a href='/'>Вернуться к форме</a></p>",
            status_code=404,
        )
    with open(json_path, "r", encoding="utf-8") as f:
        structure = json.load(f)
    slides = structure.get("slides", [])
    context = _session_context(structure, session_id)
    return render_result_page({
        "session_id": session_id,
        "title": structure.get("presentation_title", "Презентация"),
        "slide_count": len(slides),
        "preview": build_preview(slides),
        "warning": " ".join(x for x in [structure.get("generation_warning", ""), structure.get("review_warning", ""), structure.get("image_generation_warning", "")] if x).strip(),
        "source": structure.get("generation_source", ""),
        "review_source": structure.get("review_source", ""),
        "quality_review": structure.get("quality_review", {}),
        "style": normalize_style(context.get("style", DEFAULT_STYLE)),
    })


@app.get("/", response_class=HTMLResponse)
async def frontend():
    return HTML_TEMPLATE

# ── Веб-интерфейс ─────────────────────────────────────────────────────────────

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="ru">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>PresentAI — Генератор презентаций</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@400;700;800&family=DM+Sans:ital,opsz,wght@0,9..40,300;0,9..40,400;0,9..40,600;1,9..40,400&display=swap" rel="stylesheet">
<style>
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

:root {
  --bg: #07040f;
  --surface: #0e0820;
  --panel: #160d2e;
  --border: rgba(139,92,246,0.15);
  --border-bright: rgba(139,92,246,0.4);
  --accent-a: #8b5cf6;
  --accent-b: #ec4e99;
  --accent-c: #22d3ee;
  --text: #f0eaff;
  --muted: #8b7aaa;
  --radius: 16px;
  --font-head: 'Syne', sans-serif;
  --font-body: 'DM Sans', sans-serif;
}

html { scroll-behavior: smooth; }

body {
  background: var(--bg);
  color: var(--text);
  font-family: var(--font-body);
  min-height: 100vh;
  overflow-x: hidden;
}

/* ── Animated BG ─────────────────────────────────────── */
.bg-glow {
  position: fixed; inset: 0; pointer-events: none; z-index: 0;
  overflow: hidden;
}
.glow-orb {
  position: absolute; border-radius: 50%;
  filter: blur(120px); opacity: 0.25;
  animation: orbFloat 12s ease-in-out infinite;
}
.glow-orb.a { width: 600px; height: 600px; background: var(--accent-a); top: -200px; left: -200px; animation-delay: 0s; }
.glow-orb.b { width: 500px; height: 500px; background: var(--accent-b); top: 40%; right: -150px; animation-delay: -4s; }
.glow-orb.c { width: 400px; height: 400px; background: var(--accent-c); bottom: -100px; left: 30%; animation-delay: -8s; }

@keyframes orbFloat {
  0%, 100% { transform: translate(0, 0) scale(1); }
  33% { transform: translate(30px, -40px) scale(1.05); }
  66% { transform: translate(-20px, 30px) scale(0.95); }
}

/* ── Layout ──────────────────────────────────────────── */
.container {
  position: relative; z-index: 1;
  max-width: 960px; margin: 0 auto; padding: 0 24px;
}

/* ── Header ──────────────────────────────────────────── */
header {
  padding: 28px 0 0;
  display: flex; align-items: center; gap: 14px;
  animation: fadeDown 0.6s ease both;
}
.logo-mark {
  width: 48px; height: 48px; border-radius: 14px;
  background: linear-gradient(135deg, var(--accent-a), var(--accent-b));
  display: flex; align-items: center; justify-content: center;
  font-size: 22px;
  box-shadow: 0 0 24px rgba(139,92,246,0.4);
}
.logo-text { font-family: var(--font-head); font-weight: 800; font-size: 1.5rem; }
.logo-text span { color: var(--accent-a); }
.badge {
  margin-left: auto;
  font-size: 11px; font-weight: 600; letter-spacing: 0.08em; text-transform: uppercase;
  padding: 4px 10px; border-radius: 100px;
  border: 1px solid var(--border-bright);
  color: var(--accent-a);
}

/* ── Hero ────────────────────────────────────────────── */
.hero {
  padding: 60px 0 40px;
  animation: fadeUp 0.7s 0.1s ease both;
}
.hero h1 {
  font-family: var(--font-head); font-size: clamp(2rem, 5vw, 3.4rem); font-weight: 800;
  line-height: 1.1; letter-spacing: -0.02em;
  background: linear-gradient(135deg, #fff 30%, var(--accent-a), var(--accent-b));
  -webkit-background-clip: text; -webkit-text-fill-color: transparent;
  background-clip: text;
  margin-bottom: 16px;
}
.hero p {
  font-size: 1.1rem; color: var(--muted); max-width: 580px; line-height: 1.65;
}

/* ── Capability status ──────────────────────────────── */
.capability-grid {
  display: grid; grid-template-columns: repeat(4, 1fr); gap: 10px;
  margin-bottom: 30px; animation: fadeUp 0.7s 0.15s ease both;
}
.capability-item {
  min-height: 76px; padding: 12px; border-radius: 12px;
  border: 1px solid var(--border);
  background: rgba(22,13,46,0.78);
  display: grid; align-content: center; gap: 5px;
}
.capability-item strong {
  color: var(--text); font-size: 13px; font-weight: 800;
}
.capability-item span {
  color: var(--muted); font-size: 11px; line-height: 1.35;
}
.capability-item.ready { border-color: rgba(34,211,238,0.35); }
.capability-item.ready strong { color: var(--accent-c); }
@media (max-width: 760px) {
  .capability-grid { grid-template-columns: repeat(2, 1fr); }
}

/* ── Cards ───────────────────────────────────────────── */
.card {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: var(--radius);
  padding: 28px;
  animation: fadeUp 0.7s 0.2s ease both;
  transition: border-color 0.3s;
}
.card:hover { border-color: var(--border-bright); }
.card + .card { margin-top: 20px; }

.card-title {
  font-family: var(--font-head); font-weight: 700; font-size: 1rem;
  color: var(--text); margin-bottom: 18px;
  display: flex; align-items: center; gap: 8px;
}
.card-title .icon {
  width: 30px; height: 30px; border-radius: 8px;
  background: rgba(139,92,246,0.15);
  display: flex; align-items: center; justify-content: center;
  font-size: 14px;
}

/* ── Form elements ──────────────────────────────────── */
label { display: block; font-size: 13px; font-weight: 600; color: var(--muted); margin-bottom: 6px; }

textarea, input[type="text"], input[type="number"], select {
  width: 100%; padding: 12px 14px;
  background: var(--panel); border: 1px solid var(--border);
  border-radius: 10px; color: var(--text);
  font-family: var(--font-body); font-size: 15px;
  transition: border-color 0.2s, box-shadow 0.2s;
  outline: none; resize: vertical;
  -webkit-appearance: none;
}
textarea:focus, input:focus, select:focus {
  border-color: var(--accent-a);
  box-shadow: 0 0 0 3px rgba(139,92,246,0.15);
}
textarea { min-height: 100px; }
select { cursor: pointer; }
select option { background: var(--panel); }

.field { margin-bottom: 16px; }

/* Row of fields */
.fields-row { display: grid; grid-template-columns: 1fr 1fr; gap: 14px; }
@media (max-width: 600px) { .fields-row { grid-template-columns: 1fr; } }

/* ── Upload zone ─────────────────────────────────────── */
.upload-zone {
  border: 2px dashed var(--border-bright);
  border-radius: 12px; padding: 28px 20px;
  text-align: center; cursor: pointer;
  transition: all 0.2s; position: relative;
}
.upload-zone:hover, .upload-zone.dragover {
  border-color: var(--accent-a);
  background: rgba(139,92,246,0.06);
}
.upload-zone.has-file {
  border-color: var(--accent-c);
  background: rgba(34,211,238,0.07);
}
.doc-file-input {
  position: absolute;
  width: 1px; height: 1px;
  opacity: 0;
  pointer-events: none;
}
.upload-icon { font-size: 28px; margin-bottom: 8px; }
.upload-text { font-size: 14px; color: var(--muted); }
.upload-text strong { color: var(--accent-a); }
.file-name { margin-top: 10px; font-size: 13px; color: var(--accent-c); font-weight: 600; }
.file-name.invalid { color: var(--danger); }

/* ── Toggle options ──────────────────────────────────── */
.option-grid {
  display: grid; gap: 8px;
  grid-template-columns: repeat(auto-fill, minmax(120px, 1fr));
}
.option-item input { display: none; }
.option-item label {
  display: flex; flex-direction: column; align-items: center;
  gap: 6px; padding: 12px 8px; border-radius: 10px;
  border: 1px solid var(--border); cursor: pointer;
  transition: all 0.2s; color: var(--muted);
  font-size: 12px; font-weight: 600;
  text-align: center; margin: 0;
}
.option-item label .emoji { font-size: 20px; }
.option-item input:checked + label {
  border-color: var(--accent-a); background: rgba(139,92,246,0.12);
  color: var(--accent-a);
  box-shadow: 0 0 12px rgba(139,92,246,0.2);
}

/* ── Slider ──────────────────────────────────────────── */
.slider-wrap {
  --slider-progress: 25%;
  display: grid; grid-template-columns: minmax(0, 1fr) 96px; align-items: center; gap: 16px;
  width: 100%; max-width: none;
  padding: 10px 12px; border: 1px solid var(--border); border-radius: 12px;
  background: rgba(139,92,246,0.06);
}
input[type="range"] {
  width: 100%; -webkit-appearance: none; appearance: none;
  background: transparent; height: 28px; outline: none; cursor: pointer;
}
input[type="range"]:focus {
  box-shadow: none;
}
input[type="range"]::-webkit-slider-runnable-track {
  background: linear-gradient(90deg,
    var(--accent-a) 0%,
    var(--accent-b) var(--slider-progress),
    rgba(139,92,246,0.18) var(--slider-progress),
    rgba(139,92,246,0.18) 100%);
  height: 8px; border-radius: 999px;
  box-shadow: inset 0 0 0 1px rgba(255,255,255,0.12);
}
input[type="range"]::-moz-range-track {
  background: rgba(139,92,246,0.18);
  height: 8px; border-radius: 999px;
}
input[type="range"]::-moz-range-progress {
  background: linear-gradient(90deg, var(--accent-a), var(--accent-b));
  height: 8px; border-radius: 999px;
}
input[type="range"]::-webkit-slider-thumb {
  -webkit-appearance: none;
  width: 20px; height: 20px; border-radius: 50%;
  background: #fff; margin-top: -6px; cursor: pointer;
  border: 3px solid var(--accent-a);
  box-shadow: 0 0 12px rgba(139,92,246,0.55);
}
input[type="range"]::-moz-range-thumb {
  width: 20px; height: 20px; border-radius: 50%;
  background: #fff; cursor: pointer;
  border: 3px solid var(--accent-a);
  box-shadow: 0 0 12px rgba(139,92,246,0.55);
}
.slider-val {
  min-width: 112px; min-height: 54px; padding: 6px 8px;
  display: grid; grid-template-columns: 24px 1fr 24px; grid-template-rows: 1fr auto;
  align-items: center; justify-items: center; text-align: center; column-gap: 4px;
  border-radius: 12px; color: #fff;
  background: linear-gradient(135deg, rgba(139,92,246,0.95), rgba(236,78,153,0.82));
  border: 1px solid rgba(255,255,255,0.18);
  box-shadow: 0 8px 22px rgba(139,92,246,0.22);
}
.slider-val .num {
  grid-column: 2; font-feature-settings: "tnum";
  font-family: Arial, system-ui, sans-serif; font-weight: 800;
  font-size: 1.45rem; line-height: 1; color: #fff;
  font-stretch: normal; letter-spacing: 0;
  min-width: 2ch;
}
.slider-step {
  width: 24px; height: 24px; border-radius: 8px;
  border: 1px solid rgba(255,255,255,0.2);
  background: rgba(255,255,255,0.1); color: #fff;
  font-size: 16px; line-height: 1; cursor: pointer;
  display: flex; align-items: center; justify-content: center;
  padding: 0;
}
.slider-step:hover { background: rgba(255,255,255,0.18); }
.slider-step.minus { grid-column: 1; }
.slider-step.plus { grid-column: 3; }
.slider-hidden-input { display: none; }
.slider-val:focus-within {
  box-shadow: 0 0 0 3px rgba(139,92,246,0.18), 0 8px 22px rgba(139,92,246,0.22);
}
.slider-val .unit {
  grid-column: 1 / -1; display: block; margin-top: 2px; font-size: 10px;
  color: rgba(255,255,255,0.74); font-weight: 700;
}
@media (max-width: 520px) {
  .slider-wrap { width: 100%; grid-template-columns: 1fr 104px; gap: 10px; }
  input[type="range"] { width: 100%; }
}

/* ── Toggle switch ──────────────────────────────────── */
.toggle-wrap {
  display: flex; align-items: center; gap: 12px; cursor: pointer;
  padding: 12px; border: 1px solid var(--border-bright); border-radius: 12px;
  background: rgba(139,92,246,0.08);
}
.toggle { position: relative; width: 52px; height: 30px; flex: 0 0 auto; }
.toggle input {
  position: absolute; inset: 0; width: 100%; height: 100%;
  opacity: 0; cursor: pointer; z-index: 3;
}
.toggle-track {
  position: absolute; inset: 0; border-radius: 12px;
  background: var(--panel); border: 1px solid var(--border);
  transition: all 0.3s;
}
.toggle-thumb {
  position: absolute; top: 4px; left: 4px;
  width: 22px; height: 22px; border-radius: 50%;
  background: var(--muted); transition: all 0.3s;
}
.toggle input:checked ~ .toggle-track {
  background: rgba(139,92,246,0.25); border-color: var(--accent-a);
}
.toggle input:checked ~ .toggle-thumb {
  transform: translateX(22px); background: #fff;
  box-shadow: 0 0 10px rgba(139,92,246,0.5);
}
.toggle-label { font-size: 14px; color: var(--text); font-weight: 700; }

/* ── Tone cards ─────────────────────────────────────── */
.tone-grid {
  display: grid; grid-template-columns: repeat(5, 1fr); gap: 8px;
}
@media (max-width: 640px) { .tone-grid { grid-template-columns: repeat(3, 1fr); } }
.tone-card input { display: none; }
.tone-card label {
  display: flex; flex-direction: column; align-items: center; gap: 4px;
  padding: 12px 6px; border-radius: 12px; border: 1px solid var(--border);
  cursor: pointer; transition: all 0.2s; text-align: center;
  background: var(--panel); margin: 0;
}
.tone-card label:hover { border-color: var(--border-bright); background: rgba(139,92,246,0.06); }
.tone-card input:checked + label {
  border-color: var(--accent-a); background: rgba(139,92,246,0.14);
  box-shadow: 0 0 14px rgba(139,92,246,0.22);
}
.tone-icon { font-size: 22px; }
.tone-name { font-size: 12px; font-weight: 700; color: var(--text); }
.tone-desc { font-size: 10px; color: var(--muted); line-height: 1.3; }

/* ── Theme cards ─────────────────────────────────────── */
.theme-grid {
  display: grid; grid-template-columns: repeat(3, 1fr); gap: 10px;
}
@media (max-width: 600px) { .theme-grid { grid-template-columns: repeat(2, 1fr); } }
.theme-card input { display: none; }
.theme-card > label {
  display: flex; flex-direction: column; border-radius: 12px;
  border: 1px solid var(--border); overflow: hidden;
  cursor: pointer; transition: all 0.2s; margin: 0;
}
.theme-card > label:hover {
  border-color: var(--border-bright); transform: translateY(-2px);
  box-shadow: 0 6px 20px rgba(0,0,0,0.25);
}
.theme-card input:checked + label {
  border-color: var(--accent-a);
  box-shadow: 0 0 0 2px rgba(139,92,246,0.35), 0 6px 20px rgba(0,0,0,0.3);
}
.theme-preview { height: 56px; width: 100%; flex-shrink: 0; }
.theme-card-info { padding: 9px 12px; background: var(--panel); }
.theme-card-name { font-size: 13px; font-weight: 700; color: var(--text); }
.theme-card-desc { font-size: 10px; color: var(--muted); margin-top: 2px; }

/* ── Button ──────────────────────────────────────────── */
.btn-generate {
  width: 100%; margin-top: 24px; padding: 16px;
  background: linear-gradient(135deg, var(--accent-a), var(--accent-b));
  border: none; border-radius: 12px; cursor: pointer;
  font-family: var(--font-head); font-size: 1rem; font-weight: 700;
  color: white; letter-spacing: 0.02em;
  transition: all 0.2s;
  box-shadow: 0 4px 24px rgba(139,92,246,0.35);
  display: flex; align-items: center; justify-content: center; gap: 10px;
  animation: fadeUp 0.7s 0.3s ease both;
}
.btn-generate:hover { transform: translateY(-2px); box-shadow: 0 8px 32px rgba(139,92,246,0.5); }
.btn-generate:active { transform: translateY(0); }
.btn-generate:disabled { opacity: 0.5; cursor: not-allowed; transform: none; }

/* ── Progress ────────────────────────────────────────── */
#progress-section { display: none; animation: fadeUp 0.5s ease both; margin-top: 24px; }
.progress-card {
  background: var(--surface); border: 1px solid var(--border);
  border-radius: var(--radius); padding: 28px;
  text-align: center;
}
.loading-mode .hero,
.loading-mode .capability-grid,
.loading-mode #generator-form,
.loading-mode footer {
  display: none;
}
.loading-mode #progress-section {
  display: block;
  min-height: calc(100vh - 120px);
  margin-top: 34px;
}
.loading-mode .progress-card {
  min-height: 560px;
  display: flex; flex-direction: column; justify-content: center;
  border-color: var(--border-bright);
  box-shadow: 0 18px 80px rgba(139,92,246,0.18);
}
.progress-emoji { font-size: 40px; margin-bottom: 12px; animation: pulse 1.5s ease-in-out infinite; }
@keyframes pulse { 0%,100%{transform:scale(1)} 50%{transform:scale(1.1)} }
.progress-title { font-family: var(--font-head); font-size: 1.2rem; font-weight: 700; margin-bottom: 6px; }
.loading-mode .progress-title { font-size: clamp(1.6rem, 4vw, 2.45rem); }
.progress-sub { font-size: 13px; color: var(--muted); margin-bottom: 20px; }
.progress-percent {
  font-family: var(--font-head); font-size: 2.2rem; font-weight: 800;
  color: var(--accent-a); margin-bottom: 14px;
}
.progress-bar-wrap {
  background: var(--panel); border-radius: 100px; height: 8px; overflow: hidden;
}
.loading-mode .progress-bar-wrap { height: 12px; max-width: 620px; width: 100%; margin: 0 auto; }
.progress-bar {
  height: 100%; border-radius: 100px;
  background: linear-gradient(90deg, var(--accent-a), var(--accent-b));
  width: 0%; transition: width 0.5s ease;
  box-shadow: 0 0 12px rgba(139,92,246,0.5);
}
.progress-steps { margin-top: 16px; display: flex; flex-direction: column; gap: 6px; }
.progress-step { font-size: 13px; color: var(--muted); transition: color 0.3s; display: flex; align-items: center; gap: 8px; }
.progress-step.done { color: var(--text); }
.progress-step .dot { width: 6px; height: 6px; border-radius: 50%; background: var(--muted); flex-shrink: 0; transition: background 0.3s; }
.progress-step.done .dot { background: var(--accent-a); box-shadow: 0 0 6px var(--accent-a); }

/* ── Results ─────────────────────────────────────────── */
#result-section { display: none; margin-top: 24px; animation: fadeUp 0.5s ease both; }

.result-header {
  display: flex; align-items: center; gap: 14px;
  margin-bottom: 20px;
}
.result-title { font-family: var(--font-head); font-weight: 800; font-size: 1.3rem; }
.result-meta { font-size: 13px; color: var(--muted); }
.btn-download {
  margin-left: auto;
  padding: 12px 22px;
  background: linear-gradient(135deg, var(--accent-a), var(--accent-b));
  border: none; border-radius: 10px; cursor: pointer;
  font-family: var(--font-head); font-size: 0.9rem; font-weight: 700;
  color: white; transition: all 0.2s;
  box-shadow: 0 4px 16px rgba(139,92,246,0.3);
  display: flex; align-items: center; gap: 8px; text-decoration: none;
  white-space: nowrap;
}
.btn-download:hover { transform: translateY(-2px); box-shadow: 0 6px 24px rgba(139,92,246,0.45); }

/* ── Slide preview grid ─────────────────────────────── */
.slides-grid {
  display: grid; grid-template-columns: repeat(auto-fill, minmax(220px, 1fr)); gap: 14px;
}

.slide-card {
  background: var(--panel); border: 1px solid var(--border);
  border-radius: 12px; overflow: hidden;
  transition: all 0.2s;
}
.slide-card:hover { border-color: var(--border-bright); transform: translateY(-3px); box-shadow: 0 8px 24px rgba(0,0,0,0.3); }

.slide-editor { padding: 12px; display: grid; gap: 8px; }
.slide-editor input, .slide-editor textarea {
  width: 100%; min-height: 0; padding: 9px 10px; border-radius: 8px;
  font-size: 12px; color: var(--text); background: var(--surface);
  border: 1px solid var(--border);
}
.slide-editor textarea { min-height: 78px; resize: vertical; }
.btn-secondary {
  display: inline-flex; align-items: center; justify-content: center; gap: 8px;
  padding: 13px 18px; border: 1px solid var(--border-bright); border-radius: 12px;
  background: var(--surface); color: var(--text); font-weight: 700; cursor: pointer;
  transition: all 0.25s;
}
.btn-secondary:hover { border-color: var(--accent-a); transform: translateY(-2px); }
.warning-card {
  display: none; margin: 0 0 16px; padding: 12px 14px; border-radius: 12px;
  background: rgba(249,115,22,0.12); border: 1px solid rgba(249,115,22,0.35);
  color: #fed7aa; font-size: 13px; line-height: 1.45;
}
.source-badge {
  display: inline-flex; align-items: center; gap: 8px; margin: 0 0 16px;
  padding: 8px 12px; border-radius: 999px; border: 1px solid var(--border-bright);
  background: rgba(139,92,246,0.12); color: var(--text); font-size: 13px; font-weight: 700;
}
.preview-panel {
  background: var(--surface); border: 1px solid var(--border); border-radius: 14px;
  padding: 16px; margin-bottom: 18px;
}
.preview-toolbar {
  display: flex; align-items: center; justify-content: space-between; gap: 12px;
  margin-bottom: 12px; color: var(--muted); font-size: 13px;
}
.preview-slide {
  aspect-ratio: 16 / 9; border-radius: 10px; overflow: hidden;
  background: #101828; border: 1px solid rgba(255,255,255,0.08);
  padding: 5%; display: flex; flex-direction: column; justify-content: center;
  box-shadow: inset 0 0 80px rgba(119,0,255,0.16);
}
.preview-slide h2 {
  font-family: var(--font-head); font-size: clamp(1.35rem, 4vw, 2.6rem);
  line-height: 1.08; margin-bottom: 18px; color: #fff;
}
.preview-slide p, .preview-slide li { color: rgba(255,255,255,0.82); font-size: clamp(0.9rem, 2vw, 1.2rem); line-height: 1.45; }
.preview-slide ul { padding-left: 22px; display: grid; gap: 8px; }
.preview-with-image { display:grid; grid-template-columns:minmax(0,1fr) 42%; gap:18px; align-items:center; width:100%; }
.preview-image { width:100%; max-height:100%; aspect-ratio:4/3; object-fit:contain; border-radius:12px; background:rgba(255,255,255,0.06); border:1px solid rgba(255,255,255,0.12); box-shadow:0 16px 38px rgba(0,0,0,0.22); }
.preview-image-only { display:grid; gap:14px; align-items:center; justify-items:center; }
.preview-columns { display: grid; grid-template-columns: 1fr 1fr; gap: 18px; }
.preview-column { background: rgba(255,255,255,0.06); border-radius: 10px; padding: 14px; }
.preview-column h3 { color: #fff; margin-bottom: 10px; font-size: 1rem; }
.preview-stats { display: grid; grid-template-columns: repeat(3, 1fr); gap: 12px; margin-bottom: 16px; }
.preview-stat { background: rgba(255,255,255,0.08); border-radius: 10px; padding: 14px; text-align: center; }
.preview-stat strong { display: block; color: var(--accent-a); font-size: 1.35rem; margin-bottom: 4px; }
.slide-card.active { border-color: var(--accent-a); box-shadow: 0 0 0 2px rgba(139,92,246,0.22); }

.slide-thumb {
  position: relative;
  background: #0D1117;
  padding: 14px; min-height: 120px;
  display: flex; flex-direction: column; gap: 6px;
}
.slide-thumb.layout-title { background: linear-gradient(135deg, #120238, #2D1065); }
.slide-thumb.layout-conclusion { background: linear-gradient(135deg, #0D1117, #1a0533); }
.slide-thumb.layout-section_break { background: linear-gradient(135deg, #0e1a0a, #1a3510); }
.slide-thumb.layout-quote { background: linear-gradient(135deg, #0a1222, #1a2540); }
.slide-thumb.layout-stats { background: #0D1117; }
.slide-thumb.layout-two_column { background: #0D1117; }

.slide-num {
  position: absolute; top: 8px; right: 10px;
  font-size: 10px; font-weight: 700; color: rgba(255,255,255,0.3);
  font-family: var(--font-head);
}

.slide-layout-badge {
  font-size: 9px; font-weight: 700; text-transform: uppercase; letter-spacing: 0.08em;
  color: var(--accent-a); opacity: 0.8; margin-bottom: 2px;
}

.slide-thumb-title {
  font-size: 12px; font-weight: 700; color: white; line-height: 1.3;
  font-family: var(--font-head);
  display: -webkit-box; -webkit-line-clamp: 2; -webkit-box-orient: vertical; overflow: hidden;
}
.slide-thumb-content {
  font-size: 9px; color: rgba(255,255,255,0.45); line-height: 1.4;
  display: -webkit-box; -webkit-line-clamp: 3; -webkit-box-orient: vertical; overflow: hidden;
}

.slide-meta {
  padding: 8px 14px; font-size: 10px; color: var(--muted);
  display: flex; align-items: center; justify-content: space-between;
  border-top: 1px solid var(--border);
}

/* mini bullets in preview */
.mini-bullets { display: flex; flex-direction: column; gap: 3px; margin-top: 4px; }
.mini-bullet { font-size: 8px; color: rgba(255,255,255,0.45); display: flex; gap: 4px; align-items: flex-start; }
.mini-bullet::before { content: '▸'; color: var(--accent-a); flex-shrink: 0; }

/* stats preview */
.mini-stats { display: flex; gap: 6px; margin-top: 4px; flex-wrap: wrap; }
.mini-stat { background: rgba(139,92,246,0.15); border-radius: 4px; padding: 2px 6px; }
.mini-stat-val { font-size: 11px; font-weight: 700; color: var(--accent-a); }
.mini-stat-lab { font-size: 7px; color: var(--muted); }

/* ── Error ───────────────────────────────────────────── */
.error-card {
  background: rgba(239,68,68,0.08); border: 1px solid rgba(239,68,68,0.3);
  border-radius: 12px; padding: 20px; margin-top: 20px;
  display: none; color: #fca5a5; font-size: 14px;
}

/* ── RT API section ──────────────────────────────────── */
.rt-section {
  border-top: 1px solid var(--border); margin-top: 16px; padding-top: 16px;
  transition: all 0.3s;
}
.rt-section.hidden { display: none; }

/* ── Animations ──────────────────────────────────────── */
@keyframes fadeUp {
  from { opacity: 0; transform: translateY(20px); }
  to { opacity: 1; transform: translateY(0); }
}
@keyframes fadeDown {
  from { opacity: 0; transform: translateY(-10px); }
  to { opacity: 1; transform: translateY(0); }
}

/* ── Footer ──────────────────────────────────────────── */
footer {
  text-align: center; padding: 48px 24px 32px;
  font-size: 12px; color: var(--muted);
}
footer strong { color: var(--accent-a); }
</style>
</head>
<body>

<div class="bg-glow">
  <div class="glow-orb a"></div>
  <div class="glow-orb b"></div>
  <div class="glow-orb c"></div>
</div>

<div class="container">
  <header>
    <div class="logo-mark">✦</div>
    <div class="logo-text">Present<span>AI</span></div>
    <div class="badge">Амурский Код 2026</div>
  </header>

  <div class="hero">
    <h1>От промпта до<br>готового PPTX</h1>
    <p>Введите запрос, загрузите документ — и получите профессиональную презентацию за секунды. Используем LLM и AI-генерацию изображений.</p>
  </div>

  <div class="capability-grid" aria-label="Возможности сервиса">
    <div class="capability-item ready">
      <strong>RT LLM</strong>
      <span>структура и текст слайдов</span>
    </div>
    <div class="capability-item ready">
      <strong>PDF / DOCX</strong>
      <span>анализ загруженного документа</span>
    </div>
    <div class="capability-item ready">
      <strong>SD / Yandex ART</strong>
      <span>изображения по промптам</span>
    </div>
    <div class="capability-item ready">
      <strong>PPTX</strong>
      <span>сборка и скачивание файла</span>
    </div>
  </div>

  <form id="generator-form" action="/loading" method="post" enctype="multipart/form-data">
  <!-- Step 1: Prompt + Document -->
  <div class="card" id="form-section">
    <div class="card-title"><div class="icon">💬</div> Запрос и документ</div>
    <div class="field">
      <label>Описание презентации *</label>
      <textarea id="prompt" name="prompt" required placeholder="Например: Создай презентацию о преимуществах AI в телекоммуникациях для руководителей компании. Включи статистику, кейсы и план внедрения."></textarea>
    </div>
    <div class="field">
      <label>Исходный документ (необязательно)</label>
      <input class="doc-file-input" type="file" id="doc-file" name="document" accept=".pdf,.docx" onchange="return handleDocumentInputChange(this)">
      <label class="upload-zone" id="upload-zone" for="doc-file"
             ondragenter="return handleDocumentDragEnter(event)"
             ondragover="return handleDocumentDragOver(event)"
             ondragleave="return handleDocumentDragLeave(event)"
             ondrop="return handleDocumentDrop(event)">
        <div class="upload-icon">📄</div>
        <div class="upload-text">
          <strong>Выберите файл</strong> или перетащите сюда<br>
          PDF или DOCX — контент будет учтён при генерации
        </div>
        <div class="file-name" id="file-name-display" style="display:none"></div>
      </label>
      <script>
      (function () {
        window.__selectedDocumentFile = null;

        function getInput() {
          return document.getElementById('doc-file');
        }

        function getZone() {
          return document.getElementById('upload-zone');
        }

        function getDisplay() {
          return document.getElementById('file-name-display');
        }

        function isAllowed(file) {
          return !!file && /\\.(pdf|docx)$/i.test(file.name || '');
        }

        function showFile(file, errorMessage) {
          var display = getDisplay();
          var zone = getZone();
          if (!display || !zone) return;

          display.classList.toggle('invalid', !!errorMessage);
          if (errorMessage) {
            display.textContent = errorMessage;
            display.style.display = 'block';
            zone.classList.remove('has-file');
            return;
          }

          if (file) {
            display.textContent = '📎 ' + file.name;
            display.style.display = 'block';
            zone.classList.add('has-file');
          } else {
            display.textContent = '';
            display.style.display = 'none';
            zone.classList.remove('has-file');
          }
        }

        function putDroppedFileIntoInput(file) {
          var input = getInput();
          if (!input || !window.DataTransfer) return;
          try {
            var transfer = new DataTransfer();
            transfer.items.add(file);
            input.files = transfer.files;
          } catch (err) {
            // Некоторые браузеры не разрешают программно менять input.files.
          }
        }

        window.getSelectedDocumentFile = function () {
          var input = getInput();
          return window.__selectedDocumentFile || (input && input.files && input.files[0]) || null;
        };

        window.handleDocumentInputChange = function (input) {
          var file = input && input.files && input.files[0] ? input.files[0] : null;
          if (file && !isAllowed(file)) {
            input.value = '';
            window.__selectedDocumentFile = null;
            showFile(null, 'Можно загрузить только PDF или DOCX');
            return false;
          }

          window.__selectedDocumentFile = file;
          showFile(file, '');
          return true;
        };

        window.handleDocumentDragEnter = function (event) {
          if (event) event.preventDefault();
          getZone().classList.add('dragover');
          return false;
        };

        window.handleDocumentDragOver = function (event) {
          if (event) {
            event.preventDefault();
            if (event.dataTransfer) event.dataTransfer.dropEffect = 'copy';
          }
          getZone().classList.add('dragover');
          return false;
        };

        window.handleDocumentDragLeave = function (event) {
          var zone = getZone();
          if (!event || !zone.contains(event.relatedTarget)) zone.classList.remove('dragover');
          return false;
        };

        window.handleDocumentDrop = function (event) {
          var zone = getZone();
          if (event) {
            event.preventDefault();
            event.stopPropagation();
          }
          zone.classList.remove('dragover');

          var files = event && event.dataTransfer && event.dataTransfer.files;
          var file = files && files.length ? files[0] : null;
          var input = getInput();
          if (!file) return false;

          if (!isAllowed(file)) {
            if (input) input.value = '';
            window.__selectedDocumentFile = null;
            showFile(null, 'Можно загрузить только PDF или DOCX');
            return false;
          }

          window.__selectedDocumentFile = file;
          putDroppedFileIntoInput(file);
          showFile(file, '');
          return false;
        };
      })();
      </script>
    </div>

    <!-- Step 2: Settings -->
    <div class="card-title" style="margin-top:24px"><div class="icon">⚙️</div> Параметры</div>

    <div class="field">
      <label>Количество слайдов</label>
      <div class="slider-wrap">
        <input type="range" id="slide-count-range" min="4" max="20" value="8"
               data-slide-count-range>
        <div class="slider-val" id="slide-count-val">
          <button class="slider-step minus" type="button" data-slide-delta="-1" aria-label="Уменьшить количество слайдов">−</button>
          <span class="num" id="slide-count-display">8</span>
          <button class="slider-step plus" type="button" data-slide-delta="1" aria-label="Увеличить количество слайдов">+</button>
          <input class="slider-hidden-input" type="hidden" id="slide-count" name="slide_count" value="8">
          <span class="unit">слайдов</span>
        </div>
      </div>
      <script>
      (function () {
        var range = document.getElementById('slide-count-range');
        var hidden = document.getElementById('slide-count');
        var display = document.getElementById('slide-count-display');
        var unit = document.querySelector('#slide-count-val .unit');
        var minus = document.querySelector('[data-slide-delta="-1"]');
        var plus = document.querySelector('[data-slide-delta="1"]');

        function word(n) {
          var last = n % 10;
          var lastTwo = n % 100;
          if (last === 1 && lastTwo !== 11) return 'слайд';
          if ([2, 3, 4].indexOf(last) !== -1 && [12, 13, 14].indexOf(lastTwo) === -1) return 'слайда';
          return 'слайдов';
        }

        function clamp(value) {
          value = parseInt(value, 10);
          if (isNaN(value)) value = 8;
          return Math.max(4, Math.min(20, value));
        }

        function render(value) {
          value = clamp(value);
          var min = parseInt(range.min, 10);
          var max = parseInt(range.max, 10);
          var pct = ((value - min) / (max - min)) * 100;
          range.value = String(value);
          hidden.value = String(value);
          display.textContent = String(value);
          unit.textContent = word(value);
          range.parentElement.style.setProperty('--slider-progress', pct + '%');
        }

        range.addEventListener('input', function () { render(range.value); });
        range.addEventListener('change', function () { render(range.value); });
        minus.addEventListener('click', function () { render(clamp(hidden.value) - 1); });
        plus.addEventListener('click', function () { render(clamp(hidden.value) + 1); });
        render(hidden.value);
      })();
      </script>
    </div>

    <div class="field">
      <label>Тон подачи</label>
      <div class="tone-grid">
        <div class="tone-card">
          <input type="radio" name="tone" id="t-professional" value="professional" checked>
          <label for="t-professional"><span class="tone-icon">👔</span><span class="tone-name">Деловой</span><span class="tone-desc">строго и чётко</span></label>
        </div>
        <div class="tone-card">
          <input type="radio" name="tone" id="t-creative" value="creative">
          <label for="t-creative"><span class="tone-icon">🎨</span><span class="tone-name">Творческий</span><span class="tone-desc">образно и ярко</span></label>
        </div>
        <div class="tone-card">
          <input type="radio" name="tone" id="t-academic" value="academic">
          <label for="t-academic"><span class="tone-icon">🎓</span><span class="tone-name">Академический</span><span class="tone-desc">точно и строго</span></label>
        </div>
        <div class="tone-card">
          <input type="radio" name="tone" id="t-casual" value="casual">
          <label for="t-casual"><span class="tone-icon">💬</span><span class="tone-name">Дружелюбный</span><span class="tone-desc">просто и понятно</span></label>
        </div>
        <div class="tone-card">
          <input type="radio" name="tone" id="t-persuasive" value="persuasive">
          <label for="t-persuasive"><span class="tone-icon">⚡</span><span class="tone-name">Убедительный</span><span class="tone-desc">мощно и динамично</span></label>
        </div>
      </div>
    </div>

      <label>Визуальная тема</label>
      <div class="theme-grid">
        <div class="theme-card">
          <input type="radio" name="style" id="s-deep-neon" value="deep_neon" checked>
          <label for="s-deep-neon">
            <div class="theme-preview" style="background:linear-gradient(160deg,#0A0F1E 0%,#1A0A3E 100%);border-top:4px solid #7700FF"></div>
            <div class="theme-card-info"><div class="theme-card-name">Глубокий неон</div><div class="theme-card-desc">Тёмный · фиолетовый акцент</div></div>
          </label>
        </div>
        <div class="theme-card">
          <input type="radio" name="style" id="s-modern" value="modern">
          <label for="s-modern">
            <div class="theme-preview" style="background:linear-gradient(160deg,#0F172A 0%,#1E3A5F 100%);border-top:4px solid #38BDF8"></div>
            <div class="theme-card-info"><div class="theme-card-name">Тёмная волна</div><div class="theme-card-desc">Тёмный · голубой акцент</div></div>
          </label>
        </div>
        <div class="theme-card">
          <input type="radio" name="style" id="s-corporate" value="corporate">
          <label for="s-corporate">
            <div class="theme-preview" style="background:linear-gradient(160deg,#F1F5F9 0%,#DBEAFE 100%);border-top:4px solid #2563EB"></div>
            <div class="theme-card-info"><div class="theme-card-name">Деловая</div><div class="theme-card-desc">Светлый · синий акцент</div></div>
          </label>
        </div>
        <div class="theme-card">
          <input type="radio" name="style" id="s-tech" value="tech">
          <label for="s-tech">
            <div class="theme-preview" style="background:linear-gradient(160deg,#071A0E 0%,#0D2B18 100%);border-top:4px solid #22C55E"></div>
            <div class="theme-card-info"><div class="theme-card-name">Технологичная</div><div class="theme-card-desc">Тёмный · зелёный акцент</div></div>
          </label>
        </div>
        <div class="theme-card">
          <input type="radio" name="style" id="s-minimal" value="minimal">
          <label for="s-minimal">
            <div class="theme-preview" style="background:linear-gradient(160deg,#FFFFFF 0%,#F1F5F9 100%);border-top:4px solid #6366F1"></div>
            <div class="theme-card-info"><div class="theme-card-name">Минимализм</div><div class="theme-card-desc">Светлый · чистый акцент</div></div>
          </label>
        </div>
        <div class="theme-card">
          <input type="radio" name="style" id="s-creative" value="creative">
          <label for="s-creative">
            <div class="theme-preview" style="background:linear-gradient(160deg,#1E0A35 0%,#3D0F52 100%);border-top:4px solid #F97316"></div>
            <div class="theme-card-info"><div class="theme-card-name">Креативная</div><div class="theme-card-desc">Тёмный · оранжевый акцент</div></div>
          </label>
        </div>
      </div>
    </div>

    <!-- Image generation toggle -->
    <div class="rt-section">
      <div class="field">
        <label>API Token Ростелеком для LLM</label>
        <input type="text" id="rt-token" name="rt_token" placeholder="Bearer-токен для ai.rt.ru">
      </div>
      <label class="toggle-wrap" for="gen-images">
        <div class="toggle">
          <input type="checkbox" id="gen-images" name="generate_images" value="true">
          <div class="toggle-track"></div>
          <div class="toggle-thumb"></div>
        </div>
        <span class="toggle-label">Генерировать изображения (RT API)</span>
      </label>
      <div id="rt-fields" style="margin-top:14px">
        <div class="fields-row">
          <div class="field">
            <label>Модель изображений</label>
            <select id="rt-service" name="rt_service">
              <option value="yaArt">Yandex ART</option>
              <option value="sd">Stable Diffusion</option>
            </select>
          </div>
        </div>
      </div>
    </div>
  </div>

  <!-- Generate button -->
  <button class="btn-generate" id="gen-btn" type="submit">
    <span>✦</span> Сгенерировать презентацию
  </button>
  <a class="btn-secondary" style="width:100%; margin:12px 0 0; text-decoration:none" href="/api/demo">
    Скачать демо PPTX без API
  </a>
  </form>

  <!-- Progress -->
  <div id="progress-section">
    <div class="progress-card">
      <div class="progress-emoji" id="progress-emoji">🧠</div>
      <div class="progress-title" id="progress-title">Анализируем запрос...</div>
      <div class="progress-sub" id="progress-sub">LLM обрабатывает ваш запрос</div>
      <div class="progress-percent" id="progress-percent">0%</div>
      <div class="progress-bar-wrap">
        <div class="progress-bar" id="progress-bar"></div>
      </div>
      <div class="progress-steps">
        <div class="progress-step" id="ps1"><div class="dot"></div>Чтение документа</div>
        <div class="progress-step" id="ps2"><div class="dot"></div>Генерация структуры слайдов</div>
        <div class="progress-step" id="ps3"><div class="dot"></div>QA-проверка и правка слайдов</div>
        <div class="progress-step" id="ps4"><div class="dot"></div>Сборка PPTX файла</div>
      </div>
    </div>
  </div>

  <!-- Error -->
  <div class="error-card" id="error-card"></div>

  <!-- Results -->
  <div id="result-section">
    <div class="result-header">
      <div>
        <div class="result-title" id="result-title">Готово! 🎉</div>
        <div class="result-meta" id="result-meta"></div>
      </div>
      <a class="btn-download" id="download-btn" href="#" download="presentation.pptx">
        ⬇ Скачать PPTX
      </a>
      <button class="btn-secondary" id="save-edits-btn" onclick="saveEdits()">
        ↻ Обновить PPTX
      </button>
    </div>

    <div class="warning-card" id="warning-card"></div>
    <div class="source-badge" id="source-badge">Источник: ожидает генерации</div>
    <div class="preview-panel">
      <div class="preview-toolbar">
        <span id="preview-counter">Слайд 1</span>
        <span>Кликните карточку ниже, чтобы открыть слайд</span>
      </div>
      <div class="preview-slide" id="preview-slide"></div>
    </div>
    <div class="slides-grid" id="slides-grid"></div>

    <button class="btn-generate" style="margin-top:24px" onclick="resetForm()">
      ↩ Создать новую
    </button>
  </div>

</div>

<footer>
  Разработано для хакатона <strong>Амурский Код 2026</strong> · Кейс Ростелеком
</footer>

<script>
let sessionId = null;
let currentSlides = [];
let currentStyle = 'deep_neon';

window.addEventListener('error', event => {
  const message = event.message || 'Неизвестная ошибка JavaScript';
  const card = document.getElementById('error-card');
  if (card) {
    const safeMessage = String(message).replace(/[&<>"']/g, ch => ({
      '&': '&amp;', '<': '&lt;', '>': '&gt;', '"': '&quot;', "'": '&#39;'
    }[ch]));
    card.style.display = 'block';
    card.innerHTML = `<strong>Ошибка JavaScript:</strong> ${safeMessage}`;
  }
});

function toggleImages() {
  const cb = document.getElementById('gen-images');
  cb.checked = !cb.checked;
  document.getElementById('rt-fields').classList.toggle('hidden', !cb.checked);
}

function setStep(n) {
  // Старый визуальный stepper заменён карточками возможностей; функция оставлена для совместимости.
}

function setProgress(pct, emoji, title, sub) {
  document.getElementById('progress-bar').style.width = pct + '%';
  document.getElementById('progress-percent').textContent = Math.round(pct) + '%';
  if (emoji) document.getElementById('progress-emoji').textContent = emoji;
  if (title) document.getElementById('progress-title').textContent = title;
  if (sub) document.getElementById('progress-sub').textContent = sub;
}

function doneStep(id) {
  document.getElementById(id).classList.add('done');
}

function clampSlideCount(value) {
  return Math.max(4, Math.min(20, Number(value) || 8));
}

function updateSlideCountUI(value) {
  const slider = document.getElementById('slide-count-range');
  const hidden = document.getElementById('slide-count');
  const display = document.getElementById('slide-count-display');
  const unit = document.querySelector('#slide-count-val .unit');
  value = clampSlideCount(value);
  slider.value = String(value);
  hidden.value = String(value);
  display.textContent = String(value);
  unit.textContent = getSlideWord(value);
  const min = Number(slider.min);
  const max = Number(slider.max);
  const pct = ((value - min) / (max - min)) * 100;
  slider.closest('.slider-wrap').style.setProperty('--slider-progress', `${pct}%`);
}

function syncSlideCountFromRange(slider) {
  updateSlideCountUI(slider.value);
}

function syncSlideCountFromNumber(input) {
  updateSlideCountUI(input.value);
}

function changeSlideCount(delta) {
  const current = Number(document.getElementById('slide-count').value);
  updateSlideCountUI(current + delta);
}

function bindSlideCountControls() {
  const slider = document.getElementById('slide-count-range');
  const hidden = document.getElementById('slide-count');
  slider.addEventListener('input', () => updateSlideCountUI(slider.value));
  slider.addEventListener('change', () => updateSlideCountUI(slider.value));
  document.addEventListener('click', event => {
    const button = event.target.closest('[data-slide-delta]');
    if (!button) return;
    event.preventDefault();
    updateSlideCountUI(Number(hidden.value) + Number(button.dataset.slideDelta));
  });
  updateSlideCountUI(hidden.value);
}

function getSlideWord(value) {
  const last = value % 10;
  const lastTwo = value % 100;
  if (last === 1 && lastTwo !== 11) return 'слайд';
  if ([2, 3, 4].includes(last) && ![12, 13, 14].includes(lastTwo)) return 'слайда';
  return 'слайдов';
}

function showError(message) {
  document.body.classList.remove('loading-mode');
  document.getElementById('progress-section').style.display = 'none';
  const errCard = document.getElementById('error-card');
  errCard.style.display = 'block';
  errCard.innerHTML = `<strong>Ошибка:</strong> ${escapeHtml(message)}`;
  const btn = document.getElementById('gen-btn');
  btn.disabled = false;
  btn.innerHTML = '<span>✦</span> Сгенерировать презентацию';
  setStep(1);
}

async function generate() {
  let progressTimer = null;
  try {
    const prompt = document.getElementById('prompt').value.trim();
    if (!prompt) {
      showError('Введите описание презентации');
      return;
    }

    const genImages = document.getElementById('gen-images').checked;
    const rtToken = document.getElementById('rt-token').value.trim();
    const rtService = document.getElementById('rt-service').value;
    const slideCount = document.getElementById('slide-count').value;
    const style = document.querySelector('input[name="style"]:checked')?.value || 'deep_neon';
    const tone  = document.querySelector('input[name="tone"]:checked')?.value  || 'professional';
    currentStyle = style;

    const btn = document.getElementById('gen-btn');
    btn.disabled = true;
    btn.innerHTML = '<span>✦</span> Генерируем...';
    document.getElementById('error-card').style.display = 'none';
    document.getElementById('result-section').style.display = 'none';
    ['ps1','ps2','ps3','ps4'].forEach(id => document.getElementById(id).classList.remove('done'));
    document.body.classList.add('loading-mode');
    document.getElementById('progress-section').style.display = 'block';
    window.scrollTo({ top: 0, behavior: 'smooth' });
    setStep(3);
    setProgress(3, '📄', 'Генерация презентации', 'Подготавливаем запрос и документ');

    const loadingSteps = [
      { at: 12, emoji: '📄', title: 'Чтение входных данных', sub: 'Учитываем промпт, PDF/DOCX и параметры', done: 'ps1' },
      { at: 34, emoji: '🧠', title: 'RT LLM формирует структуру, QA LLM проверяет', sub: 'Создаем структуру и проверяем факты/перегруз текста', done: 'ps2' },
      { at: 62, emoji: '🎨', title: 'Готовим визуальную часть', sub: 'Обрабатываем стиль и изображения, если они включены', done: 'ps3' },
      { at: 86, emoji: '📦', title: 'Собираем PPTX', sub: 'Формируем файл презентации для скачивания', done: 'ps4' },
    ];
    let visualProgress = 3;
    progressTimer = setInterval(() => {
      visualProgress = Math.min(96, visualProgress + (visualProgress < 60 ? 3 : 1));
      const active = [...loadingSteps].reverse().find(step => visualProgress >= step.at);
      if (active) {
        doneStep(active.done);
        setProgress(visualProgress, active.emoji, active.title, active.sub);
      } else {
        setProgress(visualProgress);
      }
    }, 700);

    const formData = new FormData();
    formData.append('prompt', prompt);
    formData.append('slide_count', slideCount);
    formData.append('style', style);
    formData.append('tone', tone);
    formData.append('generate_images', genImages);
    formData.append('rt_token', rtToken);
    formData.append('rt_service', rtService);

    const fileInput = document.getElementById('doc-file');
    const documentFile = (window.getSelectedDocumentFile && window.getSelectedDocumentFile()) || fileInput.files[0];
    if (documentFile) {
      formData.append('document', documentFile);
    }

    const controller = new AbortController();
    const timeoutId = setTimeout(() => controller.abort(), 360000);
    const resp = await fetch('/api/generate', {
      method: 'POST',
      body: formData,
      signal: controller.signal
    });
    clearTimeout(timeoutId);

    const data = await resp.json().catch(() => ({}));

    if (!resp.ok) throw new Error(data.detail || 'Ошибка генерации');

    clearInterval(progressTimer);
    doneStep('ps4');
    setProgress(100, '🎉', 'Готово!', 'Презентация создана, открываем результат');

    setTimeout(() => {
      document.body.classList.remove('loading-mode');
      showResults(data);
    }, 800);

  } catch (err) {
    if (progressTimer) clearInterval(progressTimer);
    const message = err.name === 'AbortError'
      ? 'Сервер не ответил за 6 минут. Попробуйте временно выключить генерацию изображений или уменьшить количество слайдов.'
      : (err.message || String(err));
    showError(message);
  }
}

function showResults(data) {
  sessionId = data.session_id;
  currentSlides = data.preview || [];
  setStep(4);

  document.getElementById('progress-section').style.display = 'none';
  document.getElementById('result-section').style.display = 'block';
  document.getElementById('result-title').textContent = data.title || 'Презентация готова!';
  document.getElementById('result-meta').textContent = `${data.slide_count} слайдов`;
  document.getElementById('download-btn').href = `/api/download/${data.session_id}`;
  document.getElementById('warning-card').style.display = data.warning ? 'block' : 'none';
  document.getElementById('warning-card').textContent = data.warning || '';
  document.getElementById('source-badge').textContent = data.source ? `Источник: ${data.source}` : 'Источник: неизвестен';

  const grid = document.getElementById('slides-grid');
  grid.innerHTML = '';

  currentSlides.forEach(s => {
    grid.appendChild(buildSlideCard(s));
  });
  renderPreview(0);

  document.getElementById('gen-btn').disabled = false;
  document.getElementById('gen-btn').innerHTML = '<span>✦</span> Сгенерировать презентацию';
}

const LAYOUT_LABELS = {
  title: 'Титульный', content: 'Контент', two_column: 'Два столбца',
  stats: 'Статистика', quote: 'Цитата', section_break: 'Раздел',
  image: 'Изображение', conclusion: 'Заключение', section: 'Раздел',
};

function escapeHtml(value) {
  return String(value ?? '').replace(/[&<>"']/g, ch => ({
    '&': '&amp;', '<': '&lt;', '>': '&gt;', '"': '&quot;', "'": '&#39;'
  }[ch]));
}

function buildSlideCard(s) {
  const card = document.createElement('div');
  card.className = 'slide-card';
  card.dataset.index = String((s.index || 1) - 1);
  card.addEventListener('click', event => {
    if (event.target.closest('input, textarea')) return;
    renderPreview(Number(card.dataset.index));
  });

  const layout = (s.layout || 'content').replace('-', '_');
  const thumbCls = `slide-thumb layout-${layout}`;

  let contentHtml = '';

  if (s.stats && s.stats.length > 0) {
    contentHtml = `<div class="mini-stats">${s.stats.slice(0,3).map(st =>
      `<div class="mini-stat"><div class="mini-stat-val">${escapeHtml(st.value)}</div><div class="mini-stat-lab">${escapeHtml(st.label)}</div></div>`
    ).join('')}</div>`;
  } else if (s.bullets && s.bullets.length > 0) {
    contentHtml = `<div class="mini-bullets">${s.bullets.slice(0,3).map(b =>
      `<div class="mini-bullet">${escapeHtml(b.substring(0,50))}${b.length>50?'...':''}</div>`
    ).join('')}</div>`;
  } else if (s.subtitle) {
    contentHtml = `<div class="slide-thumb-content">${escapeHtml(s.subtitle)}</div>`;
  } else if (s.leftTitle || s.rightTitle) {
    contentHtml = `<div style="display:flex;gap:6px;margin-top:4px">
      <div style="flex:1;background:rgba(139,92,246,0.15);border-radius:4px;padding:4px 6px;font-size:9px;color:rgba(255,255,255,0.5)">${escapeHtml(s.leftTitle||'')}</div>
      <div style="flex:1;background:rgba(236,78,153,0.15);border-radius:4px;padding:4px 6px;font-size:9px;color:rgba(255,255,255,0.5)">${escapeHtml(s.rightTitle||'')}</div>
    </div>`;
  }

  const editorBody = (s.bullets && s.bullets.length)
    ? s.bullets.join('\n')
    : (s.content || s.subtitle || s.quote || '');

  card.innerHTML = `
    <div class="${thumbCls}">
      <div class="slide-num">${s.index}</div>
      <div class="slide-layout-badge">${LAYOUT_LABELS[layout] || layout}</div>
      <div class="slide-thumb-title">${escapeHtml(s.title || '')}</div>
      ${contentHtml}
      ${s.hasImage ? '<div style="margin-top:6px;font-size:9px;color:rgba(139,92,246,0.7)">🖼 + изображение</div>' : ''}
    </div>
    <div class="slide-editor">
      <input class="edit-title" value="${escapeHtml(s.title || '')}" placeholder="Заголовок">
      <textarea class="edit-body" placeholder="Текст, пункты или цитата">${escapeHtml(editorBody)}</textarea>
    </div>
    <div class="slide-meta">
      <span>${LAYOUT_LABELS[layout] || layout}</span>
      <span style="color:var(--accent-a)">→</span>
    </div>`;
  return card;
}

function slideContentHtml(s) {
  if (!s) return '';
  if (s.stats && s.stats.length) {
    return `<div class="preview-stats">${s.stats.slice(0, 3).map(st =>
      `<div class="preview-stat"><strong>${escapeHtml(st.value)}</strong><span>${escapeHtml(st.label)}</span></div>`
    ).join('')}</div><p>${escapeHtml(s.content || '')}</p>`;
  }
  if (s.layout === 'two_column') {
    return `<div class="preview-columns">
      <div class="preview-column"><h3>${escapeHtml(s.leftTitle || 'Слева')}</h3><ul>${(s.leftContent || []).map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul></div>
      <div class="preview-column"><h3>${escapeHtml(s.rightTitle || 'Справа')}</h3><ul>${(s.rightContent || []).map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul></div>
    </div>`;
  }
  if (s.bullets && s.bullets.length) {
    return `<ul>${s.bullets.map(x => `<li>${escapeHtml(x)}</li>`).join('')}</ul>`;
  }
  const text = s.quote || s.content || s.subtitle || '';
  return `<p>${escapeHtml(text)}</p>`;
}

function slideImageHtml(s) {
  if (!s || !s.imageData) return '';
  return `<img class="preview-image" src="${escapeHtml(s.imageData)}" alt="Изображение слайда">`;
}

function slideBodyHtml(s) {
  const body = slideContentHtml(s);
  const image = slideImageHtml(s);
  if (!image) return body;
  if (!body || body === '<p></p>') return `<div class="preview-image-only">${image}</div>`;
  return `<div class="preview-with-image"><div>${body}</div>${image}</div>`;
}
function renderPreview(index) {
  const slide = currentSlides[index];
  if (!slide) return;
  document.querySelectorAll('.slide-card').forEach((card, i) => {
    card.classList.toggle('active', i === index);
  });
  document.getElementById('preview-counter').textContent = `Слайд ${index + 1} из ${currentSlides.length}`;
  document.getElementById('preview-slide').innerHTML = `
    <h2>${escapeHtml(slide.title || `Слайд ${index + 1}`)}</h2>
    ${slideBodyHtml(slide)}
  `;
}

function collectEditedSlides() {
  return [...document.querySelectorAll('.slide-card')].map(card => {
    const idx = Number(card.dataset.index);
    const original = currentSlides[idx] || {};
    const body = card.querySelector('.edit-body').value.trim();
    const title = card.querySelector('.edit-title').value.trim();
    const slide = {...original, title};

    if (slide.layout === 'title') {
      slide.subtitle = body;
    } else if (slide.layout === 'quote') {
      slide.quote = body;
    } else if (slide.layout === 'conclusion' || slide.layout === 'section_break') {
      slide.content = body;
      slide.subtitle = body;
    } else if (slide.layout === 'stats') {
      slide.content = body;
    } else if (slide.layout === 'two_column') {
      const lines = body.split('\n').map(x => x.trim()).filter(Boolean);
      const half = Math.ceil(lines.length / 2);
      slide.leftContent = lines.slice(0, half);
      slide.rightContent = lines.slice(half);
    } else {
      slide.bullets = body.split('\n').map(x => x.trim()).filter(Boolean);
    }
    return slide;
  });
}

async function saveEdits() {
  if (!sessionId) return;
  const btn = document.getElementById('save-edits-btn');
  const slides = collectEditedSlides();
  btn.disabled = true;
  btn.textContent = 'Собираем PPTX...';
  try {
    const resp = await fetch(`/api/rebuild/${sessionId}`, {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({
        title: document.getElementById('result-title').textContent,
        style: currentStyle,
        slides
      })
    });
    const data = await resp.json();
    if (!resp.ok) throw new Error(data.detail || 'Ошибка пересборки');
    currentSlides = data.preview || slides;
    document.getElementById('download-btn').href = `/api/download/${sessionId}?t=${Date.now()}`;
    btn.textContent = 'PPTX обновлен';
    setTimeout(() => { btn.textContent = '↻ Обновить PPTX'; }, 1400);
  } catch (err) {
    alert(err.message);
    btn.textContent = '↻ Обновить PPTX';
  } finally {
    btn.disabled = false;
  }
}

function resetForm() {
  document.body.classList.remove('loading-mode');
  document.getElementById('result-section').style.display = 'none';
  document.getElementById('progress-section').style.display = 'none';
  document.getElementById('error-card').style.display = 'none';
  document.getElementById('gen-btn').disabled = false;
  document.getElementById('gen-btn').innerHTML = '<span>✦</span> Сгенерировать презентацию';
  setStep(1);
  // Сбрасываем состояние прогресса перед новой генерацией.
  ['ps1','ps2','ps3','ps4'].forEach(id => document.getElementById(id).classList.remove('done'));
  document.getElementById('progress-bar').style.width = '0%';
  document.getElementById('progress-percent').textContent = '0%';
}

// Держим currentStyle синхронным с выбранной визуальной темой.
document.querySelectorAll('input[name="style"]').forEach(r => {
  r.addEventListener('change', () => { currentStyle = r.value; });
});
currentStyle = document.querySelector('input[name="style"]:checked')?.value || 'deep_neon';

document.getElementById('generator-form').addEventListener('submit', function(event) {
  const prompt = document.getElementById('prompt').value.trim();
  if (!prompt) {
    event.preventDefault();
    showError('Введите описание презентации');
  }
});

bindSlideCountControls();

</script>
</body>
</html>
"""

# ── Точка входа ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
